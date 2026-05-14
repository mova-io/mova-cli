"""Generate a Mova iO ↔ MDK mapping slide deck (Phase J-5).

Single-slide deck (plus title + summary slides) summarising the
Mova iO platform → MDK coverage scorecard. Designed to drop into
Deva's review meeting as an embedded slide; lives standalone so it
can be regenerated as features ship.

Source of truth is ``docs/mova-io-mapping.md`` — when that doc
updates, re-run this script to regenerate the deck.

Style mirrors ``build-azure-summary-ppt.py``: Movate blue header
bar, Consolas code blocks, 16:9 widescreen. Self-contained slides
so the deck can be presented standalone or excerpted.

Usage
-----

::

    uv run --with python-pptx scripts/build-mova-io-mapping-ppt.py

Output: ``~/.movate/movate-mova-io-mapping.pptx`` by default;
override with ``--out path/to.pptx``.

No secrets, no bearer tokens — safe to email broadly.
"""

from __future__ import annotations

import argparse
import sys
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.util import Inches, Pt
except ImportError:
    print(
        "python-pptx is required. Re-run via `uv run --with python-pptx "
        "scripts/build-mova-io-mapping-ppt.py`.",
        file=sys.stderr,
    )
    sys.exit(1)


# ---------------------------------------------------------------------------
# Style constants — match other Movate decks
# ---------------------------------------------------------------------------


MOVATE_BLUE = RGBColor(0x1F, 0x3D, 0x7A)
MOVATE_GREEN = RGBColor(0x2D, 0x9C, 0x5A)
MOVATE_AMBER = RGBColor(0xD9, 0x8E, 0x04)
MOVATE_RED = RGBColor(0xC0, 0x2B, 0x2B)
TEXT_DARK = RGBColor(0x22, 0x22, 0x22)
TEXT_DIM = RGBColor(0x66, 0x66, 0x66)

SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)


def _header_bar(slide, title: str) -> None:
    """Draw the Movate-blue header bar with the slide title.

    Same shape every deck uses — keeps the brand consistent and
    saves the operator from positioning a title box manually.
    """
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, Inches(0.6))
    bar.fill.solid()
    bar.fill.fore_color.rgb = MOVATE_BLUE
    bar.line.fill.background()

    title_box = slide.shapes.add_textbox(
        Inches(0.4), Inches(0.05), SLIDE_WIDTH - Inches(0.8), Inches(0.5)
    )
    tf = title_box.text_frame
    tf.text = title
    p = tf.paragraphs[0]
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)


def _body_box(slide, top: Inches, left: Inches | None = None, height: Inches | None = None):
    """Standard left-padded body region under the header."""
    return slide.shapes.add_textbox(
        left or Inches(0.5),
        top,
        SLIDE_WIDTH - (left or Inches(0.5)) - Inches(0.5),
        height or Inches(6),
    )


# ---------------------------------------------------------------------------
# Slide builders
# ---------------------------------------------------------------------------


def _build_title_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout

    # Centered title block
    title_box = slide.shapes.add_textbox(
        Inches(1), Inches(2.5), SLIDE_WIDTH - Inches(2), Inches(1.5)
    )
    tf = title_box.text_frame
    tf.text = "Mova iO → MDK Mapping"
    p = tf.paragraphs[0]
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = MOVATE_BLUE

    subtitle = slide.shapes.add_textbox(Inches(1), Inches(4), SLIDE_WIDTH - Inches(2), Inches(0.8))
    stf = subtitle.text_frame
    stf.text = "Phase J coverage scorecard — 2026-05-14"
    sp = stf.paragraphs[0]
    sp.font.size = Pt(20)
    sp.font.color.rgb = TEXT_DIM


def _build_scorecard_slide(prs: Presentation) -> None:
    """Per-layer coverage scorecard — the headline slide.

    Renders the table from docs/mova-io-mapping.md as a PPT table.
    Each row gets a color in the Delta column so the eye lands on
    the biggest wins (Safe AI, Data & Knowledge).
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _header_bar(slide, "Coverage scorecard — before vs after Phase J")

    rows = [
        ("Layer", "Before", "After", "Delta"),
        ("AI Consumption", "85%", "85%", "—"),
        ("Agent Marketplace", "50%", "75%", "+25%"),
        ("Agent Creation & Orchestration", "80%", "95%", "+15%"),
        ("Safe AI", "40%", "80%", "+40%"),
        ("Model Layer", "70%", "80%", "+10%"),
        ("Data & Knowledge", "30%", "65%", "+35%"),
        ("AI Infrastructure", "100%", "100%", "—"),
        ("Aggregate", "~65%", "~85%", "+20%"),
    ]

    table_shape = slide.shapes.add_table(
        rows=len(rows),
        cols=4,
        left=Inches(0.5),
        top=Inches(0.9),
        width=Inches(12),
        height=Inches(4),
    )
    table = table_shape.table

    # Column widths — Layer is widest, three numeric columns narrower
    table.columns[0].width = Inches(5.5)
    for i in (1, 2, 3):
        table.columns[i].width = Inches(2.2)

    delta_column_idx = 3  # last column carries the +X% delta; colored on non-zero

    # Header row formatting
    for i, label in enumerate(rows[0]):
        cell = table.cell(0, i)
        cell.text = label
        cell.fill.solid()
        cell.fill.fore_color.rgb = MOVATE_BLUE
        for para in cell.text_frame.paragraphs:
            for run in para.runs:
                run.font.size = Pt(14)
                run.font.bold = True
                run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    # Body rows
    for r_idx, row in enumerate(rows[1:], start=1):
        for c_idx, value in enumerate(row):
            cell = table.cell(r_idx, c_idx)
            cell.text = value
            for para in cell.text_frame.paragraphs:
                for run in para.runs:
                    run.font.size = Pt(13)
                    run.font.color.rgb = TEXT_DARK
                    if r_idx == len(rows) - 1:  # Aggregate row
                        run.font.bold = True
            # Color the Delta column for non-aggregate rows
            if c_idx == delta_column_idx and row[3] not in ("—", "+0%"):
                cell.fill.solid()
                cell.fill.fore_color.rgb = MOVATE_GREEN
                for para in cell.text_frame.paragraphs:
                    for run in para.runs:
                        run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
                        run.font.bold = True

    # Caption below the table
    caption = _body_box(slide, top=Inches(5.2), height=Inches(1.5))
    ctf = caption.text_frame
    ctf.word_wrap = True
    p = ctf.paragraphs[0]
    p.text = (
        "Aggregate coverage: ~65% → ~85%. Remaining 15% concentrated in "
        "vector store (Sprint 7), memory engine (Sprint 7), and compose "
        "(Sprint 8). See docs/mova-io-mapping.md for per-box detail."
    )
    p.font.size = Pt(13)
    p.font.color.rgb = TEXT_DIM


def _build_what_shipped_slide(prs: Presentation) -> None:
    """The 14-PR roll-up — what Phase J actually delivered."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _header_bar(slide, "Phase J — what shipped (14 PRs)")

    items = [
        ("J-pre-0", "mdk add + 5 role templates", "PR #5"),
        ("J-pre-1", "project-mode validate / eval", "PR #6"),
        ("J-0", "Safe AI MVP (PII / topic / content)", "PR #8"),
        ("interrupt", "mdk list", "PR #9"),
        ("J-1", "Reflection pattern", "PR #12"),
        ("polish", "mdk guardrails CLI wrapper", "PR #13"),
        ("polish", "mdk export json-schema", "PR #14"),
        ("J-2", "mdk explain <run-id>", "PR #16"),
        ("J-3", "mdk plan --from (LLM bootstrapper)", "PR #17"),
        ("J-4", "RAG surface (knowledge.yaml)", "PR #18"),
        ("J-5", "This mapping doc + slide", "this PR"),
    ]

    body = _body_box(slide, top=Inches(0.9), height=Inches(6))
    tf = body.text_frame
    tf.word_wrap = True

    for i, (phase, feature, pr) in enumerate(items):
        para = tf.add_paragraph() if i else tf.paragraphs[0]
        para.text = f"  {phase:>12}    {feature:<40}    {pr}"
        para.font.name = "Consolas"
        para.font.size = Pt(13)
        para.font.color.rgb = TEXT_DARK

    foot = slide.shapes.add_textbox(Inches(0.5), Inches(6.4), SLIDE_WIDTH - Inches(1), Inches(0.6))
    ft = foot.text_frame
    p = ft.paragraphs[0]
    p.text = "~9,000 lines added · 1,500+ tests passing · auto-merge enabled on the repo"
    p.font.size = Pt(12)
    p.font.italic = True
    p.font.color.rgb = TEXT_DIM


def _build_demo_arc_slide(prs: Presentation) -> None:
    """The happy-path demo — what an operator does end-to-end now."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _header_bar(slide, "Demo arc — full happy path")

    body = _body_box(slide, top=Inches(0.9), height=Inches(6))
    tf = body.text_frame
    tf.word_wrap = True

    steps = [
        ("# Bootstrap from natural language", ""),
        ('$ mdk plan "Triage tickets and reply" --apply --target ./demo', ""),
        ("", ""),
        ("# Add a knowledge base", ""),
        ("$ mdk knowledge add ./docs/runbook.md --id runbook", ""),
        ('$ mdk knowledge query "SLA for P1?"', ""),
        ("", ""),
        ("# Enable safety", ""),
        ("$ mdk guardrails enable input.pii", ""),
        ('$ mdk guardrails test "leak: jane@acme.com"', ""),
        ("", ""),
        ("# Run, discover, inspect", ""),
        ("$ mdk run triage '{...}'", ""),
        ("$ mdk list", ""),
        ("$ mdk explain cccccccc", ""),
        ("", ""),
        ("# Gate + deploy", ""),
        ("$ mdk eval --project --gate 0.7", ""),
        ("$ mdk deploy", ""),
    ]

    for i, (line, _) in enumerate(steps):
        para = tf.add_paragraph() if i else tf.paragraphs[0]
        para.text = line
        para.font.name = "Consolas"
        para.font.size = Pt(13)
        if line.startswith("#"):
            para.font.color.rgb = TEXT_DIM
            para.font.italic = True
        elif line.startswith("$"):
            para.font.color.rgb = MOVATE_BLUE
            para.font.bold = True
        else:
            para.font.color.rgb = TEXT_DARK


def _build_whats_next_slide(prs: Presentation) -> None:
    """Sprint roadmap — what's queued next per BACKLOG Group L."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _header_bar(slide, "What's next — Sprint roadmap (BACKLOG Group L)")

    rows = [
        ("Sprint", "Theme", "Items"),
        ("S1", "State foundation", "snapshot / diff / rollback / audit"),
        ("S2", "Env management", "profiles / secrets / migrate / promote"),
        ("S3", "Onboarding polish", "init --project / doctor fix / openapi"),
        ("S4", "Observability polish", "monitor / tune / explain v2"),
        ("S5", "Interop & export", "oci-bundle / langgraph / simulate"),
        ("S6", "Production validation", "benchmark live / audit v2 / e2e CI"),
        ("S7", "Memory architecture", "engine + CLI (4-5 weeks)"),
        ("S8", "Multi-agent / Phase 7", "LangGraph swap-in + compose"),
    ]

    table_shape = slide.shapes.add_table(
        rows=len(rows),
        cols=3,
        left=Inches(0.5),
        top=Inches(0.9),
        width=Inches(12),
        height=Inches(5),
    )
    table = table_shape.table
    table.columns[0].width = Inches(1.0)
    table.columns[1].width = Inches(3.0)
    table.columns[2].width = Inches(8.0)

    # Header row
    for i, label in enumerate(rows[0]):
        cell = table.cell(0, i)
        cell.text = label
        cell.fill.solid()
        cell.fill.fore_color.rgb = MOVATE_BLUE
        for para in cell.text_frame.paragraphs:
            for run in para.runs:
                run.font.size = Pt(14)
                run.font.bold = True
                run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    # Body
    for r_idx, row in enumerate(rows[1:], start=1):
        for c_idx, value in enumerate(row):
            cell = table.cell(r_idx, c_idx)
            cell.text = value
            for para in cell.text_frame.paragraphs:
                for run in para.runs:
                    run.font.size = Pt(12)
                    run.font.color.rgb = TEXT_DARK

    foot = slide.shapes.add_textbox(Inches(0.5), Inches(6.4), SLIDE_WIDTH - Inches(1), Inches(0.6))
    ft = foot.text_frame
    p = ft.paragraphs[0]
    p.text = (
        "~16 weeks for CLI surface (S1-S6) + ~10 weeks for engine work (S7-S8). "
        "See BACKLOG.md Group L for dependencies + parallelisation guidance."
    )
    p.font.size = Pt(12)
    p.font.italic = True
    p.font.color.rgb = TEXT_DIM


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------


def build_deck(out: Path) -> None:
    """Assemble the 5-slide deck and write to ``out``."""
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    _build_title_slide(prs)
    _build_scorecard_slide(prs)
    _build_what_shipped_slide(prs)
    _build_demo_arc_slide(prs)
    _build_whats_next_slide(prs)

    out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(out)
    print(f"wrote {out}")


def main() -> None:
    parser = argparse.ArgumentParser(description=__doc__.split("\n")[0])
    parser.add_argument(
        "--out",
        type=Path,
        default=Path.home() / ".movate" / "movate-mova-io-mapping.pptx",
        help="Output path for the .pptx file.",
    )
    args = parser.parse_args()
    build_deck(args.out)


if __name__ == "__main__":
    main()
