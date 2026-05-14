"""Generate the Mova iO MDK Integration walkthrough deck.

Renders a 13-slide ``.pptx`` walkthrough of the MDK runtime's HTTP
surface, intended for a live screenshare with the Mova iO Angular
team. Mirrors the structure of ``docs/deva-endpoint-smoke-guide.md``
but optimized for projector-friendly reading.

Credentials never live in this script. The runtime URL + bearer
token are pulled from environment variables at run time, so the
script is safe to commit to git. The generated ``.pptx`` lands at
``~/.movate/deva-onboarding.pptx`` (outside the repo) — same
location convention as the markdown Deva-ready file.

Usage
-----

::

    export MDK_BASE="https://movate-dev-api.<hash>.eastus2.azurecontainerapps.io"
    export MDK_TOKEN="mvt_live_<your-bearer>"
    uv run --with python-pptx scripts/build-deva-ppt.py

The ``--with python-pptx`` flag installs python-pptx into a
throwaway env just for this run — no pyproject.toml pollution.
Re-run any time the values rotate; each run overwrites the output
file.

Design notes
------------

* One concept per slide. Curl + expected response side-by-side
  where it fits; otherwise stacked.
* Code blocks use Consolas at 10pt — small enough that a typical
  curl fits on one line, big enough to read from the back of a
  meeting room.
* Movate-ish blue header bar so the deck doesn't look like a
  PowerPoint default template.
* Slide numbers in the footer so the audience can reference
  "slide 7" mid-discussion.
"""

from __future__ import annotations

import os
import sys
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.util import Inches, Pt
except ImportError:
    print(
        "python-pptx not installed.\n\n"
        "Run with uv:\n"
        "  uv run --with python-pptx scripts/build-deva-ppt.py\n\n"
        "Or install into your active env:\n"
        "  pip install python-pptx",
        file=sys.stderr,
    )
    sys.exit(1)


# ---------------------------------------------------------------------------
# Inputs
# ---------------------------------------------------------------------------


def _required_env(name: str) -> str:
    value = os.environ.get(name, "").strip()
    if not value:
        print(
            f"ERROR: {name} is not set.\n\n"
            "Export both MDK_BASE and MDK_TOKEN, then re-run. Example:\n\n"
            '  export MDK_BASE="https://movate-dev-api.<hash>.eastus2.azurecontainerapps.io"\n'
            '  export MDK_TOKEN="mvt_live_<your-bearer>"\n'
            "  uv run --with python-pptx scripts/build-deva-ppt.py\n",
            file=sys.stderr,
        )
        sys.exit(2)
    return value


MDK_BASE = _required_env("MDK_BASE")
MDK_TOKEN = _required_env("MDK_TOKEN")
TENANT_ID = os.environ.get("MDK_TENANT_ID", "<your-tenant-uuid>")
CORS_ORIGINS = os.environ.get("MDK_CORS_ALLOWED_ORIGINS", "http://localhost:4200")

OUTPUT_PATH = Path.home() / ".movate" / "deva-onboarding.pptx"


# ---------------------------------------------------------------------------
# Styling
# ---------------------------------------------------------------------------

SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

MOVATE_BLUE = RGBColor(0x00, 0x4A, 0x8F)
MOVATE_BLUE_DARK = RGBColor(0x00, 0x2E, 0x5C)
TEXT_DARK = RGBColor(0x1A, 0x1A, 0x1A)
TEXT_MUTED = RGBColor(0x55, 0x55, 0x55)
CODE_BG = RGBColor(0xF5, 0xF5, 0xF5)
CODE_FG = RGBColor(0x20, 0x20, 0x20)
ACCENT_GREEN = RGBColor(0x16, 0x73, 0x4B)
ACCENT_AMBER = RGBColor(0xB5, 0x6A, 0x00)


def add_header_bar(slide, title: str, subtitle: str | None = None) -> None:
    """Movate-blue bar across the top with the slide's title."""
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, Inches(0.85)
    )
    bar.line.fill.background()
    bar.fill.solid()
    bar.fill.fore_color.rgb = MOVATE_BLUE

    title_tf = slide.shapes.add_textbox(
        Inches(0.4), Inches(0.12), Inches(12.5), Inches(0.6)
    ).text_frame
    title_tf.margin_left = title_tf.margin_right = 0
    title_p = title_tf.paragraphs[0]
    title_p.text = title
    title_p.runs[0].font.size = Pt(26)
    title_p.runs[0].font.bold = True
    title_p.runs[0].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    title_p.runs[0].font.name = "Calibri"

    if subtitle:
        sub_box = slide.shapes.add_textbox(
            Inches(0.4), Inches(0.55), Inches(12.5), Inches(0.3)
        )
        sub_p = sub_box.text_frame.paragraphs[0]
        sub_p.text = subtitle
        sub_p.runs[0].font.size = Pt(12)
        sub_p.runs[0].font.color.rgb = RGBColor(0xE0, 0xE8, 0xF4)
        sub_p.runs[0].font.name = "Calibri"


def add_text_block(
    slide,
    text: str,
    *,
    left: float,
    top: float,
    width: float,
    height: float,
    font_size: int = 14,
    color: RGBColor = TEXT_DARK,
    bold: bool = False,
) -> None:
    """Plain prose text block."""
    box = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = box.text_frame
    tf.word_wrap = True
    for i, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        if p.runs:
            run = p.runs[0]
            run.font.size = Pt(font_size)
            run.font.color.rgb = color
            run.font.bold = bold
            run.font.name = "Calibri"


def add_code_block(
    slide,
    code: str,
    *,
    left: float,
    top: float,
    width: float,
    height: float,
    font_size: int = 11,
    label: str | None = None,
) -> None:
    """Monospaced code/curl block with light-grey background."""
    if label:
        label_box = slide.shapes.add_textbox(
            Inches(left), Inches(top - 0.3), Inches(width), Inches(0.3)
        )
        lp = label_box.text_frame.paragraphs[0]
        lp.text = label
        lp.runs[0].font.size = Pt(11)
        lp.runs[0].font.bold = True
        lp.runs[0].font.color.rgb = TEXT_MUTED
        lp.runs[0].font.name = "Calibri"

    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height),
    )
    bg.line.color.rgb = RGBColor(0xCC, 0xCC, 0xCC)
    bg.fill.solid()
    bg.fill.fore_color.rgb = CODE_BG

    text_box = slide.shapes.add_textbox(
        Inches(left + 0.1),
        Inches(top + 0.05),
        Inches(width - 0.2),
        Inches(height - 0.1),
    )
    tf = text_box.text_frame
    tf.word_wrap = True
    for i, line in enumerate(code.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        if p.runs:
            run = p.runs[0]
            run.font.size = Pt(font_size)
            run.font.color.rgb = CODE_FG
            run.font.name = "Consolas"


def add_footer(slide, page_number: int, total: int) -> None:
    box = slide.shapes.add_textbox(
        Inches(0.4), Inches(7.05), Inches(12.5), Inches(0.3)
    )
    p = box.text_frame.paragraphs[0]
    p.text = (
        f"MDK + Mova iO Integration  ·  Slide {page_number} of {total}  ·  "
        "Movate Confidential"
    )
    p.runs[0].font.size = Pt(9)
    p.runs[0].font.color.rgb = TEXT_MUTED
    p.runs[0].font.name = "Calibri"


# ---------------------------------------------------------------------------
# Slide builders
# ---------------------------------------------------------------------------


def slide_title(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT
    )
    bg.line.fill.background()
    bg.fill.solid()
    bg.fill.fore_color.rgb = MOVATE_BLUE

    title_box = slide.shapes.add_textbox(
        Inches(0.6), Inches(2.6), Inches(12), Inches(1.4)
    )
    tp = title_box.text_frame.paragraphs[0]
    tp.text = "MDK + Mova iO Integration"
    tp.runs[0].font.size = Pt(44)
    tp.runs[0].font.bold = True
    tp.runs[0].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    tp.runs[0].font.name = "Calibri"

    subtitle_box = slide.shapes.add_textbox(
        Inches(0.6), Inches(4.0), Inches(12), Inches(1)
    )
    sp = subtitle_box.text_frame.paragraphs[0]
    sp.text = "Endpoint walkthrough for the Angular team"
    sp.runs[0].font.size = Pt(22)
    sp.runs[0].font.color.rgb = RGBColor(0xE0, 0xE8, 0xF4)
    sp.runs[0].font.name = "Calibri"

    foot_box = slide.shapes.add_textbox(
        Inches(0.6), Inches(6.4), Inches(12), Inches(0.4)
    )
    fp = foot_box.text_frame.paragraphs[0]
    fp.text = "Movate · MDK runtime v0.7.0 · Movate Azure Sandbox"
    fp.runs[0].font.size = Pt(13)
    fp.runs[0].font.color.rgb = RGBColor(0xB5, 0xC8, 0xE0)
    fp.runs[0].font.name = "Calibri"


def slide_connection_info(prs: Presentation, page: int, total: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header_bar(
        slide,
        "Your connection info",
        "Save these to your password manager · treat the bearer like a password",
    )
    info = (
        f"Runtime URL:   {MDK_BASE}\n"
        f"OpenAPI spec:  {MDK_BASE}/api/v1/openapi.json\n"
        f"Bearer token:  {MDK_TOKEN}\n"
        f"Tenant ID:     {TENANT_ID}\n"
        f"CORS allow:    {CORS_ORIGINS}  (your prod Mova iO origin gets added when you share the hostname)"
    )
    add_code_block(
        slide, info, left=0.4, top=1.2, width=12.5, height=2.7, font_size=14
    )
    notes = (
        "Don't commit the bearer to git.\n"
        "Don't paste it in tickets, screenshots, or public channels.\n"
        "If a previously-good request starts returning 401, the key was revoked or rotated\n"
        "    — ping the MDK team for a fresh one (~30s to mint)."
    )
    add_text_block(
        slide, "Bearer hygiene",
        left=0.4, top=4.2, width=12.5, height=0.4, font_size=18, bold=True
    )
    add_text_block(
        slide, notes,
        left=0.4, top=4.7, width=12.5, height=2.0, font_size=14, color=TEXT_DARK
    )
    add_footer(slide, page, total)


def slide_overview(prs: Presentation, page: int, total: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header_bar(slide, "What you'll exercise today",
                   "18 endpoints · 8 sections · ~15 minutes hands-on")

    inventory = (
        "1. Liveness         GET /healthz, GET /ready\n"
        "2. Browse           GET /agents, GET /api/v1/agents/{name}\n"
        "3. Create + edit    POST /api/v1/agents/from-wizard, POST /validate\n"
        "4. Run an agent     POST /runs (sync + async), GET /jobs/{id}, GET /runs/{id}, GET /trace\n"
        "5. Evaluate         POST /evals, GET /evals/{id}, GET /evals?agent=\n"
        "6. Delete           DELETE /api/v1/agents/{name}\n"
        "7. Coming soon      POST /publish, GET /history (feature-flagged today — returns 503)\n"
        "8. Troubleshooting  401 / 404 / 422 / 429 / 503 reference table"
    )
    add_code_block(
        slide, inventory, left=0.4, top=1.2, width=12.5, height=4.8, font_size=14
    )

    cta = (
        "Run the sections in order. Every curl that follows uses the runtime URL + bearer "
        "from the previous slide — copy, paste, run."
    )
    add_text_block(
        slide, cta, left=0.4, top=6.2, width=12.5, height=0.5,
        font_size=13, color=TEXT_MUTED
    )
    add_footer(slide, page, total)


def slide_endpoint(
    prs: Presentation,
    page: int,
    total: int,
    *,
    title: str,
    subtitle: str,
    purpose: str,
    curl: str,
    response: str,
    notes: str | None = None,
) -> None:
    """Generic two-column endpoint slide: curl on left, response on right."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header_bar(slide, title, subtitle)

    add_text_block(
        slide, purpose, left=0.4, top=1.05, width=12.5, height=0.6,
        font_size=14, color=TEXT_MUTED
    )

    add_code_block(
        slide, curl, left=0.4, top=2.0, width=6.3, height=4.4,
        font_size=10, label="curl"
    )
    add_code_block(
        slide, response, left=7.0, top=2.0, width=5.9, height=4.4,
        font_size=10, label="expected response"
    )

    if notes:
        add_text_block(
            slide, notes, left=0.4, top=6.5, width=12.5, height=0.5,
            font_size=12, color=ACCENT_AMBER, bold=True
        )
    add_footer(slide, page, total)


def slide_section_5_evals(prs: Presentation, page: int, total: int) -> None:
    """Evals section is dense; gets its own bespoke layout."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header_bar(
        slide, "Section 5 · Evaluate an agent",
        "Score the agent against its dataset · mean_score, pass_rate, dimensional breakdown"
    )

    add_text_block(
        slide,
        "Three endpoints work together — kick off, fetch scorecard, list history.",
        left=0.4, top=1.05, width=12.5, height=0.4, font_size=13, color=TEXT_MUTED
    )

    curls = (
        f"# 5.1 kick off (mock=true keeps it cheap)\n"
        f"curl -X POST \\\n"
        f"  \"{MDK_BASE}/api/v1/agents/faq-agent/evals\" \\\n"
        f"  -H \"Authorization: Bearer {MDK_TOKEN[:20]}...\" \\\n"
        f"  -d '{{\"gate\": 0.7, \"runs\": 1, \"mock\": true}}'\n"
        f"\n"
        f"# 5.2 fetch scorecard\n"
        f"curl \"{MDK_BASE}/api/v1/evals/<eval_id>\" \\\n"
        f"  -H \"Authorization: Bearer {MDK_TOKEN[:20]}...\"\n"
        f"\n"
        f"# 5.3 history per agent\n"
        f"curl \"{MDK_BASE}/api/v1/evals?agent=faq-agent\" \\\n"
        f"  -H \"Authorization: Bearer {MDK_TOKEN[:20]}...\""
    )
    add_code_block(
        slide, curls, left=0.4, top=1.7, width=7.5, height=4.7,
        font_size=10, label="curls"
    )

    notes = (
        "What you'll see\n"
        "─────────────────\n"
        "• 5.1 returns {eval_id, status, message}\n"
        "    status='success' → eval ran\n"
        "    status='failed' + message about missing\n"
        "    dataset → wizard agents need a dataset\n"
        "    uploaded before evals work (item 111)\n"
        "\n"
        "• 5.2 returns {mean_score, pass_rate,\n"
        "    sample_count, total_cost_usd, ...}\n"
        "\n"
        "• 5.3 returns {evals: [...], count: N}\n"
        "    sorted newest-first\n"
        "\n"
        "Render in Angular as:\n"
        "  • Score chip on agent profile\n"
        "  • Sparkline (mean_score over time)\n"
        "  • Pass/fail gate indicator"
    )
    add_code_block(
        slide, notes, left=8.2, top=1.7, width=4.7, height=4.7, font_size=11
    )
    add_footer(slide, page, total)


def slide_section_7_coming_soon(prs: Presentation, page: int, total: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header_bar(
        slide, "Section 7 · Coming soon · GitHub publish + history",
        "Routes in /openapi.json today · feature-flagged · 503 until ops registers the GitHub App"
    )

    body = (
        f"Today\n"
        f"─────\n"
        f"POST /api/v1/agents/{{name}}/publish\n"
        f"GET  /api/v1/agents/{{name}}/history\n"
        f"\n"
        f"Both return:\n"
        f"  HTTP 503\n"
        f"  {{\"detail\": {{\"error\":\n"
        f"    {{\"code\": \"agent_persistence_unavailable\",\n"
        f"     \"message\": \"github integration is disabled;\n"
        f"      set MDK_GITHUB_ENABLED=1 ...\"}}}}}}\n"
        f"\n"
        f"What it means for you\n"
        f"──────────────────────\n"
        f"The routes are already advertised in /openapi.json\n"
        f"so your `npm run client:gen` (or equivalent) picks\n"
        f"up typed methods today.\n"
        f"\n"
        f"When ops flips the flag, real responses replace the\n"
        f"503s with no Angular client changes required."
    )
    add_code_block(
        slide, body, left=0.4, top=1.2, width=7.2, height=5.4, font_size=11
    )

    future_state = (
        f"When live\n"
        f"─────────\n"
        f"\n"
        f"POST /publish returns:\n"
        f"  {{commit_sha, commit_url,\n"
        f"   branch, files_changed[]}}\n"
        f"\n"
        f"Use case: Mova iO's 'Publish'\n"
        f"button — one click commits the\n"
        f"agent bundle to the tenant's\n"
        f"GitHub repo as one atomic commit.\n"
        f"\n"
        f"GET /history returns:\n"
        f"  {{agent, commits[], page,\n"
        f"   limit, has_more}}\n"
        f"\n"
        f"Each commit: {{sha, message,\n"
        f"  author_name, author_email,\n"
        f"  timestamp, html_url}}\n"
        f"\n"
        f"Use case: version-history\n"
        f"panel on the agent profile."
    )
    add_code_block(
        slide, future_state, left=7.9, top=1.2, width=5.0, height=5.4, font_size=10
    )
    add_footer(slide, page, total)


def slide_troubleshooting(prs: Presentation, page: int, total: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header_bar(
        slide, "Troubleshooting · the 5 errors you'll actually hit",
        "Status code · cause · fix"
    )

    table = (
        "401 auth_required\n"
        "    Bearer missing, malformed, expired, or revoked.\n"
        "    Re-set MDK_TOKEN. If still 401, ping the MDK team for a fresh bearer.\n"
        "\n"
        "404 not_found\n"
        "    Agent name in the URL doesn't exist.\n"
        "    List agents (Section 2.1) to confirm spelling.\n"
        "\n"
        "422 invalid_bundle\n"
        "    Body fails schema validation.\n"
        "    Check the error.message — usually a missing required field.\n"
        "\n"
        "429 rate_limited\n"
        "    Hit per-bearer limit (60 req/min default).\n"
        "    Slow polling, or ping the MDK team for a higher limit.\n"
        "\n"
        "503 agent_persistence_unavailable\n"
        "    GitHub integration not enabled yet (publish/history endpoints only).\n"
        "    Expected — nothing to fix on your side."
    )
    add_code_block(
        slide, table, left=0.4, top=1.2, width=12.5, height=5.0, font_size=12
    )

    add_text_block(
        slide,
        "If you hit something not in this table, drop the curl + the full response "
        "(status code + body) in Slack and the MDK team will diagnose.",
        left=0.4, top=6.4, width=12.5, height=0.4,
        font_size=12, color=TEXT_MUTED
    )
    add_footer(slide, page, total)


def slide_wrap_up(prs: Presentation, page: int, total: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header_bar(
        slide, "Wrap-up · you're ready to integrate",
        "Generate the typed Angular client · ping the MDK team with questions"
    )

    next_steps = (
        f"1. Generate your typed Angular service layer\n"
        f"   ─────────────────────────────────────────\n"
        f"   Point your openapi-generator (or ng-openapi-gen) at:\n"
        f"     {MDK_BASE}/api/v1/openapi.json\n"
        f"   You get typed methods for every endpoint we walked\n"
        f"   through today.\n"
        f"\n"
        f"2. Wire the bearer\n"
        f"   ─────────────────\n"
        f"   Add an HTTP interceptor that attaches\n"
        f"     Authorization: Bearer {MDK_TOKEN[:20]}...\n"
        f"   to every outgoing request to {MDK_BASE.split('//')[-1]}.\n"
        f"\n"
        f"3. Tell us your prod hostname\n"
        f"   ──────────────────────────\n"
        f"   We'll add it to CORS via a one-line Bicep param update.\n"
        f"\n"
        f"4. Reach out\n"
        f"   ─────────\n"
        f"   Anything not in this deck — Teams DM the MDK team.\n"
        f"   Anything broken — paste the curl + full response."
    )
    add_code_block(
        slide, next_steps, left=0.4, top=1.2, width=12.5, height=5.3, font_size=12
    )
    add_footer(slide, page, total)


# ---------------------------------------------------------------------------
# Build the deck
# ---------------------------------------------------------------------------


def build_deck() -> Presentation:
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    # Slide 1 — title
    slide_title(prs)

    total = 14

    # Slide 2 — connection info
    slide_connection_info(prs, page=2, total=total)

    # Slide 3 — overview
    slide_overview(prs, page=3, total=total)

    # Slide 4 — liveness + ready
    slide_endpoint(
        prs,
        page=4,
        total=total,
        title="Section 1 · Liveness + readiness",
        subtitle="Unauthed probes · 200 means alive, 503 from /ready means a dependency is down",
        purpose=(
            "Cheap health checks. Use /healthz as a 'is the runtime up' probe "
            "on app boot; use /ready when you want to know whether the DB is reachable."
        ),
        curl=(
            f"# Liveness — no auth needed\n"
            f"curl {MDK_BASE}/healthz\n"
            f"\n"
            f"# Readiness — deep checks (DB ping)\n"
            f"curl {MDK_BASE}/ready"
        ),
        response=(
            "# /healthz\n"
            "{\n"
            '  "status": "ok",\n'
            '  "version": "0.7.0"\n'
            "}\n"
            "\n"
            "# /ready\n"
            "{\n"
            '  "status": "ready",\n'
            '  "version": "0.7.0",\n'
            '  "checks": {"storage": "ok"}\n'
            "}"
        ),
    )

    # Slide 5 — catalog browse
    slide_endpoint(
        prs,
        page=5,
        total=total,
        title="Section 2 · Browse the agent catalog",
        subtitle="List view + per-agent detail · drives the catalog screen and agent profile",
        purpose=(
            "GET /agents returns metadata only. GET /api/v1/agents/{name} returns "
            "the full profile including schemas, dataset info, and marketplace fields."
        ),
        curl=(
            f"# List every agent\n"
            f"curl {MDK_BASE}/agents \\\n"
            f"  -H \"Authorization: Bearer {MDK_TOKEN[:20]}...\"\n"
            f"\n"
            f"# Full profile for one agent\n"
            f"curl {MDK_BASE}/api/v1/agents/faq-agent \\\n"
            f"  -H \"Authorization: Bearer {MDK_TOKEN[:20]}...\""
        ),
        response=(
            "# /agents\n"
            "{\n"
            '  "agents": [\n'
            '    {"name": "faq-agent",\n'
            '     "version": "0.1.0",\n'
            '     "description": "..."},\n'
            '    ...\n'
            "  ]\n"
            "}\n"
            "\n"
            "# /api/v1/agents/{name}\n"
            "{name, version, description, model,\n"
            " input_schema, output_schema, dataset,\n"
            " role, persona, capabilities, ...}"
        ),
        notes=(
            "Note the URL difference: /agents is the catalog list; /api/v1/agents/{name} is the per-agent detail."
        ),
    )

    # Slide 6 — wizard create + validate
    slide_endpoint(
        prs,
        page=6,
        total=total,
        title="Section 3 · Create + edit (the primary Mova iO path)",
        subtitle="POST /from-wizard accepts the Onboard Agent field shape · POST /validate lints",
        purpose=(
            "POST /from-wizard is the canonical create endpoint for Mova iO. "
            "Translates wizard fields into a canonical agent bundle on the runtime."
        ),
        curl=(
            f"curl -X POST \\\n"
            f"  {MDK_BASE}/api/v1/agents/from-wizard \\\n"
            f"  -H \"Authorization: Bearer {MDK_TOKEN[:20]}...\" \\\n"
            f"  -H \"Content-Type: application/json\" \\\n"
            f"  -d '{{\n"
            f"    \"name\": \"smoke-bot\",\n"
            f"    \"agent_provider\": \"Movate\",\n"
            f"    \"agent_type\": \"Task Agent\",\n"
            f"    \"role\": \"FAQ assistant\",\n"
            f"    \"description\": \"...\",\n"
            f"    \"agent_role\": \"You are...\",\n"
            f"    \"agent_goal\": \"...\",\n"
            f"    \"agent_prompt\": \"...\",\n"
            f"    \"reference_output\": \"...\",\n"
            f"    \"ai_model\": \"gpt-4o-mini\",\n"
            f"    \"ai_foundation\": \"azure\"\n"
            f"  }}'"
        ),
        response=(
            "# 201 Created\n"
            "{\n"
            '  "name": "smoke-bot",\n'
            '  "version": "0.1.0",\n'
            '  "description": "...",\n'
            '  "agent_dir": "smoke-bot",\n'
            '  "files_persisted": [\n'
            '    "agent.yaml",\n'
            '    "prompt.md",\n'
            '    "schema/input.json",\n'
            '    "schema/output.json"\n'
            "  ]\n"
            "}\n"
            "\n"
            "# POST /validate then returns\n"
            "{passed: bool, errors, warnings,\n"
            " cost_forecast}"
        ),
        notes=(
            "agent_role is the prompt body. role is the marketplace label (chip in the catalog). Different fields."
        ),
    )

    # Slide 7 — inline run
    slide_endpoint(
        prs,
        page=7,
        total=total,
        title="Section 4a · Run an agent · inline (?wait=true)",
        subtitle="Required for wizard-created agents · sync response · latency = LLM call time",
        purpose=(
            "Synchronous run mode. API holds the connection open, executes the agent, "
            "returns the full RunView in one round-trip. Use for wizard-created agents."
        ),
        curl=(
            f"curl -X POST \\\n"
            f"  \"{MDK_BASE}/api/v1/agents/smoke-bot/runs?wait=true\" \\\n"
            f"  -H \"Authorization: Bearer {MDK_TOKEN[:20]}...\" \\\n"
            f"  -H \"Content-Type: application/json\" \\\n"
            f"  -d '{{\"input\": {{\"input\": \"What is Movate?\"}}}}'"
        ),
        response=(
            "# 200 OK\n"
            "{\n"
            '  "run_id": "...",\n'
            '  "job_id": "...",\n'
            '  "agent": "smoke-bot",\n'
            '  "provider": "openai/gpt-4o-mini",\n'
            '  "status": "success",\n'
            '  "input": {...},\n'
            '  "output": {...},\n'
            '  "metrics": {\n'
            '    "latency_ms": 2023,\n'
            '    "tokens": {"input": 31,\n'
            '               "output": 14},\n'
            '    "cost_usd": 1.3e-05\n'
            "  }\n"
            "}"
        ),
        notes=(
            "Use a 60s HTTP timeout on the client — slow LLM responses can take 20-30s."
        ),
    )

    # Slide 8 — async run + jobs polling
    slide_endpoint(
        prs,
        page=8,
        total=total,
        title="Section 4b · Run an agent · async + polling",
        subtitle="Queue-based · 202 + job_id · poll GET /jobs/{id} until terminal",
        purpose=(
            "Use for baked-in agents when the client can wait asynchronously. Returns "
            "immediately; worker picks up the job; you poll for terminal status."
        ),
        curl=(
            f"# Submit\n"
            f"curl -X POST \\\n"
            f"  {MDK_BASE}/api/v1/agents/faq-agent/runs \\\n"
            f"  -H \"Authorization: Bearer {MDK_TOKEN[:20]}...\" \\\n"
            f"  -d '{{\"input\": {{\"question\": \"...\"}}}}'\n"
            f"\n"
            f"# Poll (every 1-2 seconds)\n"
            f"curl {MDK_BASE}/jobs/<job_id> \\\n"
            f"  -H \"Authorization: Bearer {MDK_TOKEN[:20]}...\""
        ),
        response=(
            "# POST → 202\n"
            '{"job_id": "...", "status": "queued"}\n'
            "\n"
            "# GET /jobs/{id} → 200\n"
            "{\n"
            '  "job_id": "...",\n'
            '  "kind": "agent",\n'
            '  "target": "faq-agent",\n'
            '  "status": "success",\n'
            '  "result_run_id": "...",\n'
            '  "error": null,\n'
            '  "created_at": "...",\n'
            '  "completed_at": "..."\n'
            "}\n"
            "\n"
            "Status: queued | running |\n"
            "        success | error | dead_letter"
        ),
        notes=(
            "Wizard-created agents fail in async mode with unknown_agent + a hint pointing at ?wait=true. Expected; not a bug."
        ),
    )

    # Slide 9 — fetch run output + trace
    slide_endpoint(
        prs,
        page=9,
        total=total,
        title="Section 4c · Fetch run output + timeline",
        subtitle="GET /runs/{id} for output · GET /api/v1/runs/{id}/trace for full timeline",
        purpose=(
            "After a job finishes, fetch the full output via /runs/{result_run_id}. "
            "For a tree-view timeline (spans, provider metadata), use /trace."
        ),
        curl=(
            f"# Full RunView (output + metrics)\n"
            f"curl {MDK_BASE}/runs/<run_id> \\\n"
            f"  -H \"Authorization: Bearer {MDK_TOKEN[:20]}...\"\n"
            f"\n"
            f"# Trace timeline (works for workflows too)\n"
            f"curl {MDK_BASE}/api/v1/runs/<run_id>/trace \\\n"
            f"  -H \"Authorization: Bearer {MDK_TOKEN[:20]}...\""
        ),
        response=(
            "# /runs/{id}\n"
            "Same shape as the inline-run response\n"
            "from slide 7.\n"
            "\n"
            "# /api/v1/runs/{id}/trace\n"
            "{\n"
            '  "kind": "agent",  // or "workflow"\n'
            '  "run": {...},\n'
            '  "workflow": null,  // populated for\n'
            "                     // workflow runs\n"
            '  "nodes": [],\n'
            '  "total_cost_usd": 0.000036,\n'
            '  "total_latency_ms": 2972\n'
            "}"
        ),
    )

    # Slide 10 — evals
    slide_section_5_evals(prs, page=10, total=total)

    # Slide 11 — delete
    slide_endpoint(
        prs,
        page=11,
        total=total,
        title="Section 6 · Delete (soft)",
        subtitle="Move bundle to .deleted-<name>-<timestamp>/ · recoverable for ~7 days",
        purpose=(
            "Soft-delete an agent. Bundle moves to a sibling .deleted-* directory; "
            "recoverable out-of-band by the MDK team for ~7 days before final sweep."
        ),
        curl=(
            f"# Delete\n"
            f"curl -X DELETE \\\n"
            f"  {MDK_BASE}/api/v1/agents/smoke-bot \\\n"
            f"  -H \"Authorization: Bearer {MDK_TOKEN[:20]}...\"\n"
            f"\n"
            f"# Verify it's gone\n"
            f"curl {MDK_BASE}/agents \\\n"
            f"  -H \"Authorization: Bearer {MDK_TOKEN[:20]}...\""
        ),
        response=(
            "# 200 OK\n"
            "{\n"
            '  "name": "smoke-bot",\n'
            '  "deleted_dir":\n'
            '    ".deleted-smoke-bot-1778772217"\n'
            "}\n"
            "\n"
            "# /agents — smoke-bot no longer listed"
        ),
        notes=(
            "Show a confirmation dialog in your UI — soft-delete is recoverable on our side but the user shouldn't think of it as easy undo."
        ),
    )

    # Slide 12 — coming soon (GitHub)
    slide_section_7_coming_soon(prs, page=12, total=total)

    # Slide 13 — troubleshooting (intentionally penultimate so the wrap-up lands last)
    slide_troubleshooting(prs, page=13, total=total)

    # Slide 14 — wrap-up + next steps
    slide_wrap_up(prs, page=14, total=total)

    return prs


def main() -> None:
    OUTPUT_PATH.parent.mkdir(parents=True, exist_ok=True)
    prs = build_deck()

    prs.save(OUTPUT_PATH)
    # Best-effort restrict permissions; chmod 600 so a desktop scan
    # tool doesn't surface the bearer to other accounts.
    try:
        OUTPUT_PATH.chmod(0o600)
    except Exception:  # noqa: BLE001  -- chmod not critical on non-POSIX FS
        pass

    print(f"wrote deck: {OUTPUT_PATH}")
    print("open it with: open '%s'" % OUTPUT_PATH)


if __name__ == "__main__":
    main()
