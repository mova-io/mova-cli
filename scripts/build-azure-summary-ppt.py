"""Generate a slide deck summarizing the MDK Azure infrastructure.

13 slides covering the resources deployed on Movate's
``AZLABSV2.0-Sandbox(POC)`` sub — for internal leadership briefings,
ops handoffs, or "what's running on our Azure tenant" reviews.

Distinct from ``build-deva-ppt.py`` (which targets Mova iO Angular
team — endpoint walkthrough). This deck targets infra-curious
stakeholders: cost, sizing, what's deployed where, what's coming next.

Style: same Movate-blue header bar, Consolas code blocks, 16:9
widescreen. Slides are self-contained so the deck can be presented
in any order if the audience pivots to a specific topic.

Safety: no secrets baked in. Subscription ID, resource group names,
service principal client ID, and the runtime URL are all
identifier-class information that's safe to share. Unlike the
Deva onboarding deck (which has the live bearer), this deck can be
emailed broadly without rotating anything.

Usage
-----

::

    uv run --with python-pptx scripts/build-azure-summary-ppt.py

Output lands at ``~/.movate/movate-azure-summary.pptx`` by default.
Override with ``--out path/to.pptx`` if you want it elsewhere.
"""

from __future__ import annotations

import argparse
import sys
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.util import Inches, Pt
except ImportError:
    print(
        "python-pptx not installed.\n\n"
        "Run with uv:\n"
        "  uv run --with python-pptx scripts/build-azure-summary-ppt.py\n\n"
        "Or install into your active env:\n"
        "  pip install python-pptx",
        file=sys.stderr,
    )
    sys.exit(1)


# ---------------------------------------------------------------------------
# Styling
# ---------------------------------------------------------------------------

SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

MOVATE_BLUE = RGBColor(0x00, 0x4A, 0x8F)
MOVATE_BLUE_DARK = RGBColor(0x00, 0x2E, 0x5C)
TEXT_DARK = RGBColor(0x1A, 0x1A, 0x1A)
TEXT_MUTED = RGBColor(0x55, 0x55, 0x55)
CODE_BG = RGBColor(0xF5, 0xF5, 0xF5)
CODE_FG = RGBColor(0x20, 0x20, 0x20)


def add_header_bar(slide, title: str, subtitle: str | None = None) -> None:
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, Inches(0.85))
    bar.line.fill.background()
    bar.fill.solid()
    bar.fill.fore_color.rgb = MOVATE_BLUE

    title_box = slide.shapes.add_textbox(Inches(0.4), Inches(0.12), Inches(12.5), Inches(0.6))
    tp = title_box.text_frame.paragraphs[0]
    tp.text = title
    tp.runs[0].font.size = Pt(26)
    tp.runs[0].font.bold = True
    tp.runs[0].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    tp.runs[0].font.name = "Calibri"

    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.4), Inches(0.55), Inches(12.5), Inches(0.3))
        sp = sub_box.text_frame.paragraphs[0]
        sp.text = subtitle
        sp.runs[0].font.size = Pt(12)
        sp.runs[0].font.color.rgb = RGBColor(0xE0, 0xE8, 0xF4)
        sp.runs[0].font.name = "Calibri"


def add_text_block(
    slide,
    text: str,
    *,
    left: float,
    top: float,
    width: float,
    height: float,
    font_size: int = 14,
    color: RGBColor = TEXT_DARK,
    bold: bool = False,
) -> None:
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    for i, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        if p.runs:
            run = p.runs[0]
            run.font.size = Pt(font_size)
            run.font.color.rgb = color
            run.font.bold = bold
            run.font.name = "Calibri"


def add_code_block(
    slide,
    code: str,
    *,
    left: float,
    top: float,
    width: float,
    height: float,
    font_size: int = 11,
    label: str | None = None,
) -> None:
    if label:
        label_box = slide.shapes.add_textbox(
            Inches(left), Inches(top - 0.3), Inches(width), Inches(0.3)
        )
        lp = label_box.text_frame.paragraphs[0]
        lp.text = label
        lp.runs[0].font.size = Pt(11)
        lp.runs[0].font.bold = True
        lp.runs[0].font.color.rgb = TEXT_MUTED
        lp.runs[0].font.name = "Calibri"

    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height)
    )
    bg.line.color.rgb = RGBColor(0xCC, 0xCC, 0xCC)
    bg.fill.solid()
    bg.fill.fore_color.rgb = CODE_BG

    text_box = slide.shapes.add_textbox(
        Inches(left + 0.1),
        Inches(top + 0.05),
        Inches(width - 0.2),
        Inches(height - 0.1),
    )
    tf = text_box.text_frame
    tf.word_wrap = True
    for i, line in enumerate(code.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        if p.runs:
            run = p.runs[0]
            run.font.size = Pt(font_size)
            run.font.color.rgb = CODE_FG
            run.font.name = "Consolas"


def add_footer(slide, page_number: int, total: int) -> None:
    box = slide.shapes.add_textbox(Inches(0.4), Inches(7.05), Inches(12.5), Inches(0.3))
    p = box.text_frame.paragraphs[0]
    p.text = (
        f"Movate Azure · MDK Infrastructure Summary  ·  Slide {page_number} of {total}  ·  "
        "Movate Internal"
    )
    p.runs[0].font.size = Pt(9)
    p.runs[0].font.color.rgb = TEXT_MUTED
    p.runs[0].font.name = "Calibri"


# ---------------------------------------------------------------------------
# Slide builders
# ---------------------------------------------------------------------------


def slide_title(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT)
    bg.line.fill.background()
    bg.fill.solid()
    bg.fill.fore_color.rgb = MOVATE_BLUE

    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(2.4), Inches(12), Inches(1.6))
    tp = title_box.text_frame.paragraphs[0]
    tp.text = "MDK Runtime on Movate Azure"
    tp.runs[0].font.size = Pt(44)
    tp.runs[0].font.bold = True
    tp.runs[0].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    tp.runs[0].font.name = "Calibri"

    sub_box = slide.shapes.add_textbox(Inches(0.6), Inches(3.9), Inches(12), Inches(1))
    sp = sub_box.text_frame.paragraphs[0]
    sp.text = "Infrastructure summary · AZLABSV2.0-Sandbox(POC) · deployed 2026-05-14"
    sp.runs[0].font.size = Pt(20)
    sp.runs[0].font.color.rgb = RGBColor(0xE0, 0xE8, 0xF4)
    sp.runs[0].font.name = "Calibri"

    foot_box = slide.shapes.add_textbox(Inches(0.6), Inches(6.4), Inches(12), Inches(0.4))
    fp = foot_box.text_frame.paragraphs[0]
    fp.text = "10 resources · 6 role assignments · 6 KV secrets · ~$25-40/mo · Movate Internal"
    fp.runs[0].font.size = Pt(13)
    fp.runs[0].font.color.rgb = RGBColor(0xB5, 0xC8, 0xE0)
    fp.runs[0].font.name = "Calibri"


def slide_overview(prs: Presentation, page: int, total: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header_bar(
        slide,
        "At a glance",
        "Subscription · Resource group · Live URL · Resource counts",
    )

    info = (
        "Subscription:     AZLABSV2.0-Sandbox(POC)\n"
        "Sub ID:           8fab0f8f-b577-45d7-a485-ec32f73b22be\n"
        "Tenant:           Movate Technologies Pvt Ltd\n"
        "Resource group:   movate-dev-rg\n"
        "Region:           East US 2\n"
        "Deployed:         2026-05-14 (blue/green migration from a personal pay-as-you-go sub)\n"
        "\n"
        "Live runtime URL:\n"
        "  https://movate-dev-api.bluebush-9aec1e70.eastus2.azurecontainerapps.io\n"
        "\n"
        "OpenAPI spec:\n"
        "  https://movate-dev-api.bluebush-9aec1e70.eastus2.azurecontainerapps.io/api/v1/openapi.json"
    )
    add_code_block(slide, info, left=0.4, top=1.2, width=12.5, height=3.4, font_size=14)

    summary = (
        "Resource inventory\n"
        "──────────────────\n"
        "  • 10 Azure resources in movate-dev-rg\n"
        "  • 6 role assignments wiring identities to ACR + Key Vault\n"
        "  • 6 secrets in Key Vault (Postgres + LLM provider keys)\n"
        "  • 1 service principal for unattended deploys (Contributor + UAA)\n"
        "\n"
        "Steady-state cost: ~$25-40 / month infrastructure"
    )
    add_code_block(slide, summary, left=0.4, top=4.9, width=12.5, height=2.0, font_size=14)
    add_footer(slide, page, total)


def slide_architecture(prs: Presentation, page: int, total: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header_bar(
        slide,
        "Architecture",
        "External callers → CAE → API / Worker → ACR + KV + Postgres → LLM providers",
    )

    diagram = (
        "External callers\n"
        "(Mova iO Angular, curl, CI)\n"
        "        │\n"
        "        │  HTTPS + Bearer (mvt_live_...)\n"
        "        ▼\n"
        "┌─────────────────────────────────────────────────────────────────┐\n"
        "│              movate-dev-cae (Container Apps Environment)        │\n"
        "│                                                                 │\n"
        "│   ┌─────────────────────┐        ┌─────────────────────┐        │\n"
        "│   │ movate-dev-api      │        │ movate-dev-worker   │        │\n"
        "│   │ (public, :8000)     │ ──▶◀── │ (internal)          │        │\n"
        "│   │ identity: api-mi    │  queue │ identity: worker-mi │        │\n"
        "│   └────┬───────────┬────┘        └────┬───────────┬────┘        │\n"
        "└────────┼───────────┼──────────────────┼───────────┼─────────────┘\n"
        "         │           │                  │           │\n"
        "         │ ACR pull  │ secret reads     │           │ outbound HTTPS\n"
        "         ▼           ▼                  ▼           ▼\n"
        "  ┌────────────┐ ┌────────────┐ ┌────────────────┐ ┌──────────────┐\n"
        "  │ movatedev  │ │ movate-dev │ │ movate-dev-pg- │ │  openai      │\n"
        "  │ acrmvt     │ │ -kv-mvt    │ │ mvt (Postgres) │ │  anthropic   │\n"
        "  │ (ACR)      │ │ (KV)       │ │ Flex Burstable │ │  langfuse    │\n"
        "  └────────────┘ └────────────┘ └────────────────┘ └──────────────┘\n"
        "\n"
        "        All container logs + metrics ─▶ movate-dev-logs (Log Analytics)"
    )
    add_code_block(slide, diagram, left=0.3, top=1.1, width=12.7, height=5.7, font_size=9)
    add_footer(slide, page, total)


def slide_compute(prs: Presentation, page: int, total: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header_bar(
        slide, "Compute (3 resources)", "Container Apps + the environment that hosts them"
    )

    body = (
        "movate-dev-cae · Container Apps Environment\n"
        "  Hosts both Container Apps below. Provides shared networking, log\n"
        "  routing, and the consumption-billed compute pool.\n"
        "\n"
        "movate-dev-api · Container App (public ingress :8000)\n"
        "  FastAPI HTTP runtime. Serves /api/v1/*, /healthz, /ready,\n"
        "  /openapi.json. Identity: movate-dev-api-mi (UAI).\n"
        "  Scale: 1 min / 2 max replicas, 0.5 vCPU / 1.0 GiB.\n"
        "\n"
        "movate-dev-worker · Container App (internal-only)\n"
        "  Background job worker. Claims jobs from Postgres, executes\n"
        "  against LLM providers, persists RunRecords. Identity:\n"
        "  movate-dev-worker-mi (UAI).\n"
        "  Scale: 1 min / 2 max replicas, 0.5 vCPU / 1.0 GiB.\n"
        "  Trigger: queue depth 3 per replica (KEDA postgresql scaler)."
    )
    add_code_block(slide, body, left=0.4, top=1.2, width=12.5, height=5.4, font_size=12)
    add_footer(slide, page, total)


def slide_data_storage(prs: Presentation, page: int, total: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header_bar(
        slide, "Data + storage (3 resources)", "Postgres · Key Vault · Container Registry"
    )

    body = (
        "movate-dev-pg-mvt · Azure Database for PostgreSQL Flexible Server\n"
        "  SKU:     Burstable B1ms (1 vCore, 2 GB RAM)\n"
        "  Storage: 32 GB · 7-day backup retention\n"
        "  Tables:  jobs, runs, evals, api_keys, workflow_runs\n"
        "  Conn:    sslmode=require · password via KV secret\n"
        "\n"
        "movate-dev-kv-mvt · Key Vault\n"
        "  Holds 6 runtime secrets (next slide)\n"
        "  Container Apps reference via secretRef → values never appear\n"
        "  in pod env vars, deployment outputs, or ARM history\n"
        "\n"
        "movatedevacrmvt · Container Registry\n"
        "  SKU:     Basic (admin user enabled on dev only)\n"
        "  Holds:   movate:0.7.0-<sha> + movate:0.7.0-latest\n"
        "  Pull:    via UAI (AcrPull role assignment)"
    )
    add_code_block(slide, body, left=0.4, top=1.2, width=12.5, height=5.4, font_size=12)
    add_footer(slide, page, total)


def slide_identity(prs: Presentation, page: int, total: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header_bar(
        slide,
        "Identity + IAM (3 UAIs · 6 role assignments · 1 SP)",
        "User-assigned managed identities · service principal for unattended deploys",
    )

    body = (
        "User-Assigned Managed Identities (UAIs)\n"
        "────────────────────────────────────────\n"
        "  movate-dev-api-mi         → consumed by movate-dev-api\n"
        "  movate-dev-worker-mi      → consumed by movate-dev-worker\n"
        "  movate-dev-teams-bot-mi   → pre-staged for the Teams bot\n"
        "                               (not yet deployed; UAI is free)\n"
        "\n"
        "Role Assignments (6)\n"
        "─────────────────────\n"
        "  api-mi         → AcrPull              on movatedevacrmvt\n"
        "  worker-mi      → AcrPull              on movatedevacrmvt\n"
        "  teams-bot-mi   → AcrPull              on movatedevacrmvt\n"
        "  api-mi         → KV Secrets User      on movate-dev-kv-mvt\n"
        "  worker-mi      → KV Secrets User      on movate-dev-kv-mvt\n"
        "  teams-bot-mi   → KV Secrets User      on movate-dev-kv-mvt\n"
        "\n"
        "Service Principal (1) — used by deploy scripts + CI\n"
        "────────────────────────────────────────────────────\n"
        "  fe9e2bf7-e212-4c70-a153-19e7c8a98269\n"
        "    + Contributor                on subscription\n"
        "    + User Access Administrator  on subscription\n"
        "  Credentials live in ~/.movate/azure.env (chmod 600, outside repo)"
    )
    add_code_block(slide, body, left=0.4, top=1.1, width=12.5, height=5.7, font_size=11)
    add_footer(slide, page, total)


def slide_kv_secrets(prs: Presentation, page: int, total: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header_bar(
        slide,
        "Key Vault contents (6 secrets)",
        "All referenced by Container Apps via secretRef · never appear in pod env / ARM",
    )

    body = (
        "Secret name              Reader                    Source / rotation\n"
        "─────────────────────    ──────────────────────    ─────────────────────────────\n"
        "pg-admin-password        api-mi, worker-mi         openssl rand at deploy time\n"
        "                                                   Rotate quarterly\n"
        "\n"
        "pg-connection-string     api-mi, worker-mi         Composed from FQDN + password\n"
        "                                                   at deploy time\n"
        "\n"
        "openai-api-key           api-mi, worker-mi         OpenAI dashboard\n"
        "\n"
        "anthropic-api-key        api-mi, worker-mi         Anthropic console\n"
        "\n"
        "langfuse-secret-key      api-mi, worker-mi         Langfuse project settings\n"
        "\n"
        "langfuse-public-key      api-mi, worker-mi         Langfuse project settings"
    )
    add_code_block(slide, body, left=0.4, top=1.2, width=12.5, height=5.4, font_size=11)
    add_footer(slide, page, total)


def slide_observability(prs: Presentation, page: int, total: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header_bar(
        slide,
        "Observability (1 resource)",
        "Log Analytics — collects every Container App log + metric",
    )

    body = (
        "movate-dev-logs · Log Analytics Workspace\n"
        "──────────────────────────────────────────\n"
        "  Retention:  30 days (dev tier; prod default is 90)\n"
        "  Receives:   stdout/stderr + metrics from movate-dev-api and\n"
        "              movate-dev-worker, plus the CAE's platform logs\n"
        "  Query:      KQL via the Azure Portal or `az monitor log-analytics`\n"
        "\n"
        "External observability (not provisioned by this Bicep — operator\n"
        "─────────────────────  configures via KV secrets)\n"
        "\n"
        "  Langfuse  · LLM call traces (prompt/response/tokens/cost)\n"
        "              Configured via langfuse-secret-key + langfuse-public-key\n"
        "              in KV. Trace-replay UI lives in Langfuse Cloud.\n"
        "\n"
        "  OpenTelemetry · OTel spans (workflow → node → provider call)\n"
        "                  Optional dependency; flip via env var.\n"
        "\n"
        "Operator paths\n"
        "──────────────\n"
        "  az containerapp logs show -g movate-dev-rg -n movate-dev-api --tail 100\n"
        "  az containerapp logs show -g movate-dev-rg -n movate-dev-worker --tail 100"
    )
    add_code_block(slide, body, left=0.4, top=1.2, width=12.5, height=5.4, font_size=11)
    add_footer(slide, page, total)


def slide_scale_matrix(prs: Presentation, page: int, total: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header_bar(
        slide,
        "Per-environment scale matrix",
        "Today's dev tier · staging + prod defaults documented in Bicep",
    )

    body = (
        "Aspect                  Dev (current)             Prod (future)\n"
        "─────────────────────   ────────────────────────  ────────────────────────\n"
        "API replicas            1 min / 2 max             2 min / 10 max\n"
        "Worker replicas         1 min / 2 max             2 min / 20 max\n"
        "API CPU/RAM             0.5 vCPU / 1.0 GiB        1.0 vCPU / 2.0 GiB\n"
        "Worker CPU/RAM          0.5 vCPU / 1.0 GiB        1.0 vCPU / 2.0 GiB\n"
        "Postgres SKU            Burstable B1ms / 32 GB    GeneralPurpose D2ds_v5 / 64 GB\n"
        "Postgres backups        7 days                    14 days\n"
        "ACR SKU                 Basic                     Standard\n"
        "Log Analytics retention 30 days                   90 days\n"
        "Worker scale trigger    Queue depth 3/replica     Queue depth 10/replica\n"
        "\n"
        "All defaults parameterized in infra/azure/main.bicep — flip\n"
        "`param env` to `staging` or `prod` and SKUs auto-adjust."
    )
    add_code_block(slide, body, left=0.4, top=1.2, width=12.5, height=5.4, font_size=11)
    add_footer(slide, page, total)


def slide_cost(prs: Presentation, page: int, total: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header_bar(
        slide,
        "Cost — steady-state dev tier",
        "Infrastructure only · LLM call spend is pass-through",
    )

    body = (
        "Line item                                           ~ / month\n"
        "─────────────────────────────────────────────────   ─────────\n"
        "Postgres Burstable B1ms + 32 GB storage              $13\n"
        "Azure Container Registry (Basic SKU)                  $5\n"
        "Container Apps consumption (API + Worker)          $5 to $15  (traffic-dependent)\n"
        "Log Analytics (pay-as-you-go, ~1 GB/mo)               $3\n"
        "Key Vault (transactions only)                        <$1\n"
        "Container Apps Environment                            $0  (consumption-based)\n"
        "User-assigned managed identities (x3)                 $0  (free tier)\n"
        "─────────────────────────────────────────────────   ─────────\n"
        "Total infrastructure                              ~$25 to $40\n"
        "\n"
        "LLM API calls (OpenAI / Anthropic / Langfuse)        pass-through\n"
        "                                                     (billed to those\n"
        "                                                     accounts separately\n"
        "                                                     based on usage)\n"
        "\n"
        "Scaling to prod: GeneralPurpose Postgres + Standard ACR + bumped\n"
        "Container Apps replicas roughly 3-4x infra cost. Prod typical:\n"
        "$120-180/mo, still dominated by Postgres."
    )
    add_code_block(slide, body, left=0.4, top=1.2, width=12.5, height=5.4, font_size=11)
    add_footer(slide, page, total)


def slide_not_deployed(prs: Presentation, page: int, total: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header_bar(
        slide,
        "What's NOT deployed yet",
        "Deliberately scoped out of the v0.7 migration · roadmap items",
    )

    body = (
        "Teams bot Container App + Bot Service registration\n"
        "  Status:  movate-dev-teams-bot-mi UAI is pre-staged\n"
        "  Why not: not on the Friday demo path; Mova iO Angular is the focus\n"
        "  How to:  flip `enableTeamsBot = true` in main.movate.bicepparam\n"
        "\n"
        "Custom domain / vanity URL (e.g. api.mova-io.movate.com)\n"
        "  Status:  using the Azure-generated *.azurecontainerapps.io FQDN\n"
        "  Why not: adds DNS + cert complexity for zero demo value\n"
        "  How to:  add a Container Apps custom domain + Managed Certificate\n"
        "\n"
        "Staging + prod environments\n"
        "  Status:  Bicep is parameterized for all 3 envs; only dev provisioned\n"
        "  How to:  `az group create movate-staging-rg` + apply Bicep with env=staging\n"
        "\n"
        "CI/CD wired against this sub\n"
        "  Status:  today's deploys run locally via the SP\n"
        "  How to:  add a GitHub Actions workflow calling friday-demo-deploy.sh\n"
        "\n"
        "Azure Files mount for cross-pod agent bundles (BACKLOG item 109)\n"
        "  Status:  workaround in flight via ?wait=true on /runs (item 110, live)\n"
        "  Why:     wizard-created agents land on the API pod, not the worker\n"
        "  How to:  mount an Azure Files share at /home/movate/agents on both pods"
    )
    add_code_block(slide, body, left=0.4, top=1.1, width=12.5, height=5.7, font_size=10)
    add_footer(slide, page, total)


def slide_operator_quick_ref(prs: Presentation, page: int, total: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header_bar(
        slide, "Operator quick reference", "Switching to SP context · checking state · redeploying"
    )

    body = (
        "# Switch az context to this sub\n"
        "set -a; source ~/.movate/azure.env; set +a\n"
        'az login --service-principal -u "$AZURE_CLIENT_ID" \\\n'
        '  -p "$AZURE_CLIENT_SECRET" --tenant "$AZURE_TENANT_ID"\n'
        "az account set --subscription 8fab0f8f-b577-45d7-a485-ec32f73b22be\n"
        "\n"
        "# Inventory check (matches this deck)\n"
        'az resource list -g movate-dev-rg --query "[].{name:name, type:type}" -o table\n'
        'az role assignment list --all --assignee "$AZURE_CLIENT_ID" -o table\n'
        'az keyvault secret list --vault-name movate-dev-kv-mvt --query "[].name" -o tsv\n'
        "\n"
        "# Smoke the runtime (unauth)\n"
        "curl https://movate-dev-api.bluebush-9aec1e70.eastus2.azurecontainerapps.io/healthz\n"
        "\n"
        "# Tail container logs\n"
        "az containerapp logs show -g movate-dev-rg -n movate-dev-api    --tail 100\n"
        "az containerapp logs show -g movate-dev-rg -n movate-dev-worker --tail 100\n"
        "\n"
        "# Re-deploy a new image tag\n"
        "az acr build --registry movatedevacrmvt --image movate:0.7.0-<sha> \\\n"
        "  -f Dockerfile --target runtime .\n"
        "az deployment group create -g movate-dev-rg -f infra/azure/main.bicep \\\n"
        "  -p infra/azure/main.movate.bicepparam --parameters image=movate:0.7.0-<sha>"
    )
    add_code_block(slide, body, left=0.4, top=1.2, width=12.5, height=5.4, font_size=10)
    add_footer(slide, page, total)


def slide_source_of_truth(prs: Presentation, page: int, total: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header_bar(
        slide, "Source of truth · where to look next", "Bicep · runbook · architecture deep-dive"
    )

    body = (
        "Bicep templates (what to deploy)\n"
        "─────────────────────────────────\n"
        "  infra/azure/main.bicep                 — top-level orchestrator\n"
        "  infra/azure/main.movate.bicepparam    — Movate-sub parameters (gitignored)\n"
        "  infra/azure/modules/*.bicep            — per-resource modules\n"
        "                                            (ACR, KV, Postgres, CAE,\n"
        "                                            containerapp-api,\n"
        "                                            containerapp-worker)\n"
        "\n"
        "Operator-facing docs (how to operate it)\n"
        "─────────────────────────────────────────\n"
        "  docs/azure-movate-architecture.md      — comprehensive architecture\n"
        "                                            reference (this deck = the\n"
        "                                            10-minute version)\n"
        "  docs/azure-movate-migration-runbook.md — 12-step blue/green migration\n"
        "                                            playbook (status: COMPLETE)\n"
        "  docs/azure-credentials-setup.md        — SP credential storage\n"
        "                                            convention\n"
        "\n"
        "Deploy script\n"
        "──────────────\n"
        "  scripts/friday-demo-deploy.sh          — end-to-end deploy:\n"
        "                                            ACR build → Bicep apply\n"
        "                                            → smoke → mint bearer"
    )
    add_code_block(slide, body, left=0.4, top=1.2, width=12.5, height=5.4, font_size=11)
    add_footer(slide, page, total)


def slide_close(prs: Presentation, page: int, total: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT)
    bg.line.fill.background()
    bg.fill.solid()
    bg.fill.fore_color.rgb = MOVATE_BLUE_DARK

    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(2.6), Inches(12), Inches(1.4))
    tp = title_box.text_frame.paragraphs[0]
    tp.text = "Questions · feedback · next steps"
    tp.runs[0].font.size = Pt(36)
    tp.runs[0].font.bold = True
    tp.runs[0].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    tp.runs[0].font.name = "Calibri"

    body_box = slide.shapes.add_textbox(Inches(0.6), Inches(4.0), Inches(12), Inches(2.5))
    body_text = (
        "Ping the MDK team for:\n"
        "  · Access to the sub (Reader/Contributor scoped to the RG)\n"
        "  · A live walkthrough of the runtime\n"
        "  · Roadmap (Teams bot · staging · prod · custom domain)\n"
        "\n"
        "Repository: github.com/mova-io/mova-cli"
    )
    body_box.text_frame.word_wrap = True
    for i, line in enumerate(body_text.split("\n")):
        p = body_box.text_frame.paragraphs[0] if i == 0 else body_box.text_frame.add_paragraph()
        p.text = line
        if p.runs:
            p.runs[0].font.size = Pt(18)
            p.runs[0].font.color.rgb = RGBColor(0xE0, 0xE8, 0xF4)
            p.runs[0].font.name = "Calibri"


# ---------------------------------------------------------------------------
# Build the deck
# ---------------------------------------------------------------------------


def build_deck() -> Presentation:
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    total = 13

    slide_title(prs)
    slide_overview(prs, page=2, total=total)
    slide_architecture(prs, page=3, total=total)
    slide_compute(prs, page=4, total=total)
    slide_data_storage(prs, page=5, total=total)
    slide_identity(prs, page=6, total=total)
    slide_kv_secrets(prs, page=7, total=total)
    slide_observability(prs, page=8, total=total)
    slide_scale_matrix(prs, page=9, total=total)
    slide_cost(prs, page=10, total=total)
    slide_not_deployed(prs, page=11, total=total)
    slide_operator_quick_ref(prs, page=12, total=total)
    slide_source_of_truth(prs, page=13, total=total)
    slide_close(prs, page=14, total=14)

    return prs


def main() -> None:
    parser = argparse.ArgumentParser(description="Generate the Movate Azure infra summary deck.")
    parser.add_argument(
        "--out",
        type=Path,
        default=Path.home() / ".movate" / "movate-azure-summary.pptx",
        help="Output path (default: ~/.movate/movate-azure-summary.pptx)",
    )
    args = parser.parse_args()

    args.out.parent.mkdir(parents=True, exist_ok=True)
    prs = build_deck()
    prs.save(args.out)

    print(f"wrote deck: {args.out}")
    print(f"open with: open '{args.out}'")


if __name__ == "__main__":
    main()
