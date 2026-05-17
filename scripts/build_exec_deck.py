"""Build an executive-audience .pptx deck for movate-cli.

Why a script rather than a hand-built pptx: the slide content is
prose that lives in version control alongside the code it describes.
When BACKLOG.md changes, this script regenerates the deck with the
fresh content. Designers can take the generated file as a starting
point and apply branding/themes on top.

Usage:
    uv pip install python-pptx        # one-time
    uv run python scripts/build_exec_deck.py

Output: docs/movate-cli-exec-deck.pptx
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt

# -----------------------------------------------------------------------------
# Brand colors — neutral palette; designer can re-skin in PowerPoint
# -----------------------------------------------------------------------------

INK = RGBColor(0x1A, 0x1A, 0x1A)
INK_DIM = RGBColor(0x55, 0x55, 0x55)
ACCENT = RGBColor(0x00, 0x6E, 0xB8)  # deep blue
ACCENT_LITE = RGBColor(0x00, 0xA8, 0xE8)
GREEN = RGBColor(0x2E, 0x86, 0x36)
AMBER = RGBColor(0xC1, 0x7B, 0x00)
RED = RGBColor(0xB0, 0x2A, 0x2A)
BG = RGBColor(0xFF, 0xFF, 0xFF)
BG_PANEL = RGBColor(0xF6, 0xF7, 0xF9)


# -----------------------------------------------------------------------------
# Helpers — keep slide builders small + composable
# -----------------------------------------------------------------------------


def add_title_slide(prs: Presentation, title: str, subtitle: str) -> None:
    """Cover slide. Big title centered + dim subtitle."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    _add_text(
        slide, title, left=0.6, top=2.4, width=12.0, height=1.6,
        size=54, bold=True, color=INK, align="left",
    )
    _add_text(
        slide, subtitle, left=0.6, top=4.1, width=12.0, height=1.0,
        size=24, color=INK_DIM, align="left",
    )
    _add_text(
        slide,
        "Internal AI agent platform · v1.0 review",
        left=0.6, top=6.6, width=12.0, height=0.5,
        size=14, color=ACCENT, align="left",
    )


def add_section(prs: Presentation, label: str, title: str) -> None:
    """Section divider. Small kicker + big title on a colored band."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_filled_rect(slide, left=0.0, top=0.0, width=13.33, height=7.5, color=BG_PANEL)
    _add_text(
        slide, label.upper(), left=0.6, top=2.8, width=12.0, height=0.5,
        size=14, bold=True, color=ACCENT, align="left",
    )
    _add_text(
        slide, title, left=0.6, top=3.3, width=12.0, height=1.5,
        size=44, bold=True, color=INK, align="left",
    )


def add_content_slide(
    prs: Presentation,
    title: str,
    bullets: list[str],
    *,
    subtitle: str | None = None,
    footer: str | None = None,
) -> None:
    """Generic content slide. Title at top, bullets as the body, optional
    dim subtitle and one-line footer."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    # Header rule
    _add_filled_rect(slide, left=0.6, top=1.05, width=12.13, height=0.04, color=ACCENT)
    _add_text(
        slide, title, left=0.6, top=0.4, width=12.13, height=0.7,
        size=28, bold=True, color=INK, align="left",
    )
    body_top = 1.3
    if subtitle:
        _add_text(
            slide, subtitle, left=0.6, top=1.2, width=12.13, height=0.4,
            size=14, color=INK_DIM, align="left",
        )
        body_top = 1.65

    _add_bullets(
        slide,
        bullets,
        left=0.6, top=body_top, width=12.13, height=5.5,
        size=18, color=INK,
    )

    if footer:
        _add_text(
            slide, footer, left=0.6, top=6.95, width=12.13, height=0.4,
            size=11, color=INK_DIM, align="left",
        )


def add_two_column_slide(
    prs: Presentation,
    title: str,
    left_heading: str,
    left_bullets: list[str],
    right_heading: str,
    right_bullets: list[str],
) -> None:
    """Side-by-side comparison or two-list layout."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_filled_rect(slide, left=0.6, top=1.05, width=12.13, height=0.04, color=ACCENT)
    _add_text(
        slide, title, left=0.6, top=0.4, width=12.13, height=0.7,
        size=28, bold=True, color=INK, align="left",
    )
    # Left column
    _add_text(
        slide, left_heading, left=0.6, top=1.3, width=5.9, height=0.5,
        size=18, bold=True, color=ACCENT, align="left",
    )
    _add_bullets(
        slide, left_bullets, left=0.6, top=1.85, width=5.9, height=5.3,
        size=15, color=INK,
    )
    # Right column
    _add_text(
        slide, right_heading, left=6.83, top=1.3, width=5.9, height=0.5,
        size=18, bold=True, color=ACCENT, align="left",
    )
    _add_bullets(
        slide, right_bullets, left=6.83, top=1.85, width=5.9, height=5.3,
        size=15, color=INK,
    )


def add_metrics_slide(
    prs: Presentation,
    title: str,
    metrics: list[tuple[str, str]],
    *,
    subtitle: str | None = None,
) -> None:
    """Big-number tiles. ``metrics`` = list of (value, label) pairs."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_filled_rect(slide, left=0.6, top=1.05, width=12.13, height=0.04, color=ACCENT)
    _add_text(
        slide, title, left=0.6, top=0.4, width=12.13, height=0.7,
        size=28, bold=True, color=INK, align="left",
    )
    if subtitle:
        _add_text(
            slide, subtitle, left=0.6, top=1.2, width=12.13, height=0.4,
            size=14, color=INK_DIM, align="left",
        )

    # Layout tiles in a 2 x N grid
    tile_w = 2.8
    tile_h = 1.9
    gap = 0.25
    cols = 4
    start_left = 0.6
    start_top = 2.0
    for i, (value, label) in enumerate(metrics):
        row, col = divmod(i, cols)
        left = start_left + col * (tile_w + gap)
        top = start_top + row * (tile_h + gap)
        _add_filled_rect(slide, left=left, top=top, width=tile_w, height=tile_h, color=BG_PANEL)
        _add_text(
            slide, value, left=left, top=top + 0.15, width=tile_w, height=0.95,
            size=36, bold=True, color=ACCENT, align="center",
        )
        _add_text(
            slide, label, left=left, top=top + 1.15, width=tile_w, height=0.65,
            size=13, color=INK_DIM, align="center",
        )


def add_table_slide(
    prs: Presentation,
    title: str,
    headers: list[str],
    rows: list[list[str]],
    *,
    subtitle: str | None = None,
) -> None:
    """Tabular layout. Headers + N data rows."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_filled_rect(slide, left=0.6, top=1.05, width=12.13, height=0.04, color=ACCENT)
    _add_text(
        slide, title, left=0.6, top=0.4, width=12.13, height=0.7,
        size=28, bold=True, color=INK, align="left",
    )
    if subtitle:
        _add_text(
            slide, subtitle, left=0.6, top=1.2, width=12.13, height=0.4,
            size=14, color=INK_DIM, align="left",
        )

    table_top = Inches(1.7 if subtitle else 1.4)
    table = slide.shapes.add_table(
        rows=len(rows) + 1,
        cols=len(headers),
        left=Inches(0.6),
        top=table_top,
        width=Inches(12.13),
        height=Inches(0.5 + 0.55 * len(rows)),
    ).table

    # Header row
    for c, header in enumerate(headers):
        cell = table.cell(0, c)
        cell.fill.solid()
        cell.fill.fore_color.rgb = ACCENT
        tf = cell.text_frame
        tf.text = header
        p = tf.paragraphs[0]
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = BG

    # Data rows
    for r, row in enumerate(rows, start=1):
        for c, val in enumerate(row):
            cell = table.cell(r, c)
            cell.fill.solid()
            cell.fill.fore_color.rgb = BG_PANEL if r % 2 == 0 else BG
            tf = cell.text_frame
            tf.text = val
            p = tf.paragraphs[0]
            p.font.size = Pt(12)
            p.font.color.rgb = INK


# -----------------------------------------------------------------------------
# Low-level primitives
# -----------------------------------------------------------------------------


def _add_text(
    slide,
    text: str,
    *,
    left: float,
    top: float,
    width: float,
    height: float,
    size: int = 16,
    bold: bool = False,
    italic: bool = False,
    color: RGBColor = INK,
    align: str = "left",
) -> None:
    tx = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = tx.text_frame
    tf.word_wrap = True
    tf.text = text
    p = tf.paragraphs[0]
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.italic = italic
    p.font.color.rgb = color
    if align == "center":
        from pptx.enum.text import PP_ALIGN  # noqa: PLC0415
        p.alignment = PP_ALIGN.CENTER
    elif align == "right":
        from pptx.enum.text import PP_ALIGN  # noqa: PLC0415
        p.alignment = PP_ALIGN.RIGHT


def _add_bullets(
    slide,
    items: list[str],
    *,
    left: float,
    top: float,
    width: float,
    height: float,
    size: int,
    color: RGBColor,
) -> None:
    tx = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = tx.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.space_after = Pt(8)
        p.level = 0


def _add_filled_rect(
    slide, *, left: float, top: float, width: float, height: float, color: RGBColor
) -> None:
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    shape.shadow.inherit = False


# -----------------------------------------------------------------------------
# Deck construction
# -----------------------------------------------------------------------------


def build() -> Presentation:
    prs = Presentation()
    # 16:9 widescreen
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    # ──────────────────────────────────────────────────────────────
    # Slide 1 — Title
    add_title_slide(
        prs,
        title="movate-cli",
        subtitle="Internal AI agent platform — built, deployed, validated end-to-end on Azure",
    )

    # Slide 2 — Executive summary
    add_content_slide(
        prs,
        title="Executive summary",
        subtitle="What we shipped, in one slide",
        bullets=[
            "v1.0 of an internal AI agent + workflow platform. Single binary covers the full developer loop: scaffold, iterate, evaluate, benchmark, deploy.",
            "Production-deployed on Azure Container Apps + Postgres Flex + Key Vault. Validated end-to-end against a real subscription.",
            "Multi-vendor LLM support via LiteLLM (OpenAI, Anthropic, Gemini, Azure OpenAI). No vendor lock-in in the application layer.",
            "Tenant isolation enforced at every SQL boundary. Per-tenant cost ceilings, per-job model-policy enforcement, full observability via Langfuse + OpenTelemetry.",
            "Operator alerts (email / SMS / Telegram) for production jobs. CI gates on eval regression + bench cost-drift.",
            "Feature-complete for v1.0 scope. Next investments are templates, customer-onboarding gates (privacy redaction), and v1.1 workflow features.",
        ],
    )

    # Slide 3 — The problem
    add_content_slide(
        prs,
        title="The problem",
        subtitle="Why we built this",
        bullets=[
            "Movate engineers were prototyping AI agents in scattered notebooks, isolated repos, and ad-hoc Streamlit apps — no shared definition of \"what an agent is.\"",
            "Each team picked its own LLM vendor, secrets management, deployment approach, and quality bar. Knowledge didn't compose.",
            "Cost and quality were invisible. No one knew which models a customer-facing agent was paying for, or whether a prompt edit had regressed accuracy.",
            "Operational handoffs (eng → SRE → security) had no shared substrate — every project relitigated tenant isolation, secrets, deploys, RBAC.",
            "Goal: one platform that turns \"an agent\" into a typed, version-controlled, evaluated, deployable artifact — with the same ergonomics as a Python library.",
        ],
    )

    # Slide 4 — The solution
    add_content_slide(
        prs,
        title="The solution: movate-cli",
        subtitle="One CLI. One YAML contract. One deploy story.",
        bullets=[
            "Declarative `agent.yaml` (Pydantic-validated) is the contract. Schema in, schema out — typed JSON in / typed JSON out.",
            "`movate init / validate / watch / run / eval / bench / deploy / submit` covers the entire developer loop.",
            "Multi-vendor LLM support behind a single `BaseLLMProvider` Protocol. Models are config, not code.",
            "Workflows compose agents into a graph. Linear DAGs run on a homegrown runner; conditional / parallel / human-in-the-loop opt into LangGraph.",
            "Built-in observability: Langfuse + OpenTelemetry traces, cost attribution per call, drift baselines for eval and bench.",
            "Azure-native deploy: Bicep modules + GitHub Actions OIDC federated credentials. `git push release/dev` → running revision.",
        ],
    )

    # ──────────────────────────────────────────────────────────────
    # Slide 5 — Section: WHAT'S SHIPPED
    add_section(prs, "Part 1", "What's shipped in v1.0")

    # Slide 6 — Architecture overview
    add_content_slide(
        prs,
        title="System architecture",
        subtitle="Three layers: local dev loop, Azure runtime, operator alerts",
        bullets=[
            "Local developer loop: `movate` CLI → local Executor → LiteLLM → any provider. Stores everything in sqlite for offline replay.",
            "Azure deployed runtime: API + worker on Container Apps; Postgres Flex for state; Key Vault for secrets; user-assigned managed identities for ACR pull + KV reads.",
            "KEDA Postgres scaler scales the worker on queue depth (claimable jobs), not CPU. Leading indicator, not lagging.",
            "Operator alert fan-out (MultiDispatcher pattern): every terminal job fans out to email / SMS / Telegram channels — each channel decides whether to fire.",
            "Diagram with full edges in `docs/v1.0-overview.md` (renders inline in GitHub).",
        ],
        footer="See: docs/v1.0-overview.md · infra/azure/main.bicep · docs/azure-bootstrap.md",
    )

    # Slide 7 — Local developer experience
    add_content_slide(
        prs,
        title="Shipped: Local developer experience",
        subtitle="The inner loop — from blank dir to typed, validated agent",
        bullets=[
            "`movate init <name>` — scaffolds agent.yaml, prompt.md, JSON schemas, eval dataset stub, judge config.",
            "`movate watch` — hot-reload TDD. Poll filesystem; on save, re-run validate + linter. Save the prompt, see the diff in <1s.",
            "Prompt linter (4 rules): undeclared input refs (Jinja2 AST analysis), empty prompt, missing JSON instruction, no output-schema reference. Errors exit 2; --strict promotes warnings.",
            "Cost forecast on every validate: estimates eval-run cost from dataset size × pricing-table. Catches \"$4 surprise\" bills before running.",
            "`movate run --mock` — deterministic provider, no API keys needed. Smoke-test wiring without burning provider budget.",
            "Rich + Typer panels group commands by phase. `movate --help` is the discovery surface.",
        ],
    )

    # Slide 8 — Azure deployment
    add_content_slide(
        prs,
        title="Shipped: Azure deployment",
        subtitle="One-command deploy after one-time bootstrap",
        bullets=[
            "`scripts/azure-bootstrap.sh dev` — idempotent script: resource group + service principal + federated OIDC credential + role assignments. Eliminates the most error-prone manual step.",
            "`infra/azure/main.bicep` — seven modular Bicep files (ACR, Key Vault, Postgres Flex, ACA env, two apps, etc.). Per-env SKU defaults: dev is cheap Burstable; prod is GeneralPurpose with 2 warm replicas.",
            "User-assigned managed identities created at the top level — role assignments grant access BEFORE the apps come up. Fixes a real chicken-and-egg deadlock we hit on the first cold deploy.",
            "`movate deploy --target dev` wraps `az acr build` (cloud-side image build) + `az containerapp update` (revision roll) + `/healthz` poll until version matches. Rollback via `--skip-build --image-tag <prev>`.",
            "`movate doctor --target <env>` walks the deploy path — az login → subscription → RG → ACR → both apps → /healthz. Operator pointers on every red.",
            "`.github/workflows/deploy.yml` — federated OIDC, scoped per-env GitHub Environments for approval gates. `git push release/dev` → auto-deploy.",
        ],
    )

    # Slide 9 — Quality gates
    add_content_slide(
        prs,
        title="Shipped: Quality gates as code",
        subtitle="Eval + bench with CI-gateable regression detection",
        bullets=[
            "`movate eval --gate 0.7 --runs 3` — dataset eval with exact-match or LLM-as-judge scoring. Cross-family enforcement: judge ≠ tested family (catches confounded scores).",
            "`movate eval --baseline <id> --regression-tolerance 0.05` — diffs current run vs stored EvalRecord. Exits 1 on regression past tolerance. CI-gateable.",
            "`movate bench` — multi-model comparison. Per-provider cost mean, latency p50/p95, score under the configured gate mode (mean | min | p10).",
            "`movate bench --baseline <id>` — per-model deltas vs a stored bench, with matched/added/removed model sets surfaced. Same regression gate semantics.",
            "Tenant-scoped monthly cost ceiling enforced at executor entry — zero provider cost on a budget-blocked run.",
            "Model policy (`movate.yaml: policy.allowed_providers / deny_models / max_cost_per_run_usd`) enforced at validate + at executor entry — bundles loaded by `movate serve` can't bypass.",
        ],
    )

    # Slide 10 — Observability
    add_content_slide(
        prs,
        title="Shipped: Observability + tracing",
        subtitle="Cost-aware spans on every provider call",
        bullets=[
            "Langfuse integration: every run / eval / bench produces a Langfuse trace with prompt, completion, tokens, cost, and pricing-table version stamped on the root span.",
            "OpenTelemetry alternative: OTLP exporter via env config — drops into any OTel backend (Honeycomb, Datadog, Grafana Tempo, …).",
            "Per-call span attributes: `cost_usd`, `pricing_version`, `chosen_provider`, `tokens.input / output / cached_input`. Dashboards filter on pricing_version drift without joining back to RunRecord.",
            "Cost-drift detection: every run logs both our pricing-table cost AND LiteLLM's reported cost. Drift > 5% logs a loud warning. Our table stays canonical.",
            "`movate trace replay <run-id>` — reconstructs a run / workflow's node-by-node execution from stored RunRecords. Useful for debugging weeks-old failures.",
        ],
    )

    # Slide 11 — Operator notifications
    add_two_column_slide(
        prs,
        title="Shipped: Operator notifications",
        left_heading="Per-job opt-in channels",
        left_bullets=[
            "Email (SMTP) — vendor-agnostic. Works with ACS Email, SendGrid, Mailgun, AWS SES, Gmail — anything that speaks SMTP.",
            "Activated via `--notify-email <addr>`. Worker fires fire-and-forget after each terminal transition.",
            "SMS (Azure Communication Services) — Azure-native; connection string in KV behind managed identity.",
            "Activated via `--notify-sms +1...` with client-side E.164 normalize+validate.",
            "Blocked on operator-side A2P 10DLC registration (~2-3 weeks) for US deliverability — code-ready.",
        ],
        right_heading="Operator-wide channel",
        right_bullets=[
            "Telegram bot — free, no regulatory tax, cross-platform.",
            "Operator-wide trigger: pings on EVERY terminal job (vs. per-job opt-in for email/SMS).",
            "Right shape for personal dev-loop alerts (\"my job finished\"). 5-minute setup: BotFather → /start → KV paste.",
            "Validated end-to-end against author's real phone during platform validation.",
            "MultiDispatcher composes all three: each backend decides whether the job addresses its channel.",
        ],
    )

    # Slide 12 — Tenant isolation
    add_content_slide(
        prs,
        title="Shipped: Tenant isolation as a first-class concern",
        subtitle="Every storage path filters by tenant_id at the SQL WHERE clause",
        bullets=[
            "Single multi-tenant Postgres. Every persisted row carries `tenant_id`; every read filters on it.",
            "9 audit gaps closed in v1.0 stage 4: get_run / get_workflow_run / get_eval / get_job / get_bench / update_job / revoke_api_key / touch_api_key / list_evals / list_workflow_runs.",
            "Cross-tenant lookups return `None`, never 403 — leaking 403 would leak the existence of cross-tenant IDs.",
            "Parametrized fuzz tests over all three storage backends (in-memory, sqlite, postgres) sweep every cross-tenant read path. CI gates on it.",
            "Workflow checkpointer (LangGraph integration) tenant-namespaces thread IDs — tenant A's pause/resume threads are invisible to tenant B even if IDs collide.",
            "API key tenant scope is encoded in the key prefix (`mvt_<env>_<tenant>_<keyid>_<secret>`) and validated server-side on every request.",
        ],
    )

    # Slide 13 — Workflow engine
    add_content_slide(
        prs,
        title="Shipped: Workflow engine",
        subtitle="Compose agents into typed graphs · dual-runtime · checkpointable",
        bullets=[
            "Workflow YAML compiles to an internal IR (`WorkflowGraph` with `NodeType` + `EdgeKind` enums) — runtime-agnostic.",
            "Linear DAGs run on a homegrown runner (zero deps). Conditional / parallel / human-in-the-loop opt into LangGraph via `runtime: langgraph`.",
            "LangGraph compiler supports `interrupt_before` for HUMAN nodes + tenant-namespaced checkpointers (memory / sqlite / postgres) for pause-resume across processes.",
            "Resume API: HUMAN node interrupts the graph → external system POSTs `/workflows/{id}/resume` with operator decision → execution continues from the checkpoint with the merged state.",
            "Determinism bundle (v1.1): tenant-namespaced checkpointer, resume API, HITL nodes, conditional + parallel edges, tool registry — all shipped, all CI-green.",
            "v1.1+: TOOL / FUNCTION / SUB_WORKFLOW node types currently rejected by `can_compile` with operator-facing error pointers. Each follow-up PR flips one rejection.",
        ],
    )

    # ──────────────────────────────────────────────────────────────
    # Slide 14 — Section: DEVELOPER EXPERIENCE
    add_section(prs, "Part 2", "Developer experience")

    # Slide 15 — Developer workflow (Day 0 → Day 3+)
    add_two_column_slide(
        prs,
        title="From `movate init` to running on Azure",
        left_heading="Day 0 — Scaffold + iterate",
        left_bullets=[
            "`movate init my-agent` (scaffold)",
            "`movate watch ./my-agent` (hot-reload TDD)",
            "`movate validate` (linter + cost forecast)",
            "`movate run --mock` (no API call)",
            "",
            "Day 1 — Quality gates",
            "`movate eval --gate 0.7` (regression gate)",
            "`movate eval --baseline <id>` (drift)",
            "`movate bench -m gpt-4o -m claude-3-5-haiku`",
            "`movate bench --baseline <id>` (cost drift)",
        ],
        right_heading="Day 2 — Deploy to Azure",
        right_bullets=[
            "`scripts/azure-bootstrap.sh dev`",
            "`az deployment group create ...`",
            "`movate config add-target dev ...`",
            "`movate doctor --target dev`",
            "`movate deploy --target dev`",
            "",
            "Day 3+ — Use the deployed runtime",
            "`movate submit faq-agent '...' --target dev`",
            "`movate jobs show <id>` / `jobs list`",
            "📱 Telegram bot pings on terminal",
        ],
    )

    # Slide 16 — Live demo placeholder
    add_content_slide(
        prs,
        title="Live demo (optional)",
        subtitle="5-minute walkthrough — replace this slide or skip",
        bullets=[
            "Open a terminal: `movate init faq-agent-demo`",
            "Edit `prompt.md` in your editor while `movate watch ./faq-agent-demo` runs — show the live re-validate loop with cost forecast.",
            "`movate run ./faq-agent-demo \"what is movate?\" --mock` — show the typed JSON output.",
            "`movate eval ./faq-agent-demo --gate 0.7 --mock` — show pass/fail with the Rich table.",
            "(Already-deployed dev env) `movate submit faq-agent-demo '...' --target dev --wait` — wait ~3s, Telegram phone dings on the slide.",
            "Total runtime: ~5 minutes. Pre-stage the `--target dev` setup since the deployment is live.",
        ],
        footer="Tip: pre-record a 90-second screencast as fallback in case live demo gets stuck",
    )

    # ──────────────────────────────────────────────────────────────
    # Slide 17 — Section: STATUS + COST
    add_section(prs, "Part 3", "Status, cost, risk")

    # Slide 18 — By the numbers
    add_metrics_slide(
        prs,
        title="By the numbers",
        subtitle="Snapshot at v1.0 — every number CI-verified",
        metrics=[
            ("770+", "unit tests"),
            ("0", "failing tests"),
            ("3", "storage backends parametrized"),
            ("5", "PRs in active stack"),
            ("22", "movate sub-commands"),
            ("9/9", "doctor checks green on live Azure"),
            ("~$50", "/mo idle cost on dev"),
            ("100%", "tenant-scoped queries"),
        ],
    )

    # Slide 19 — Cost analysis
    add_table_slide(
        prs,
        title="Cost analysis (Azure)",
        subtitle="Steady-state idle cost per environment — scales with traffic, not idle",
        headers=["Env", "Monthly idle (USD)", "Postgres tier", "ACA replicas", "Notes"],
        rows=[
            ["dev", "~$50", "Burstable B1ms", "1 / 1", "Sandbox; tear down when not in use"],
            ["staging", "~$100-200", "Burstable B2s", "1 / 2", "Pre-release validation"],
            ["prod", "~$300-500", "Standard_D2ds_v5", "2 / 10", "Always-warm replicas; longer log retention"],
            ["Per-environment", "+ token costs", "(passthrough)", "(passthrough)", "LLM usage billed by provider; pricing-table tracks"],
        ],
    )

    # Slide 20 — Security + risk posture
    add_content_slide(
        prs,
        title="Security + risk posture",
        subtitle="Where v1.0 stands and what's covered",
        bullets=[
            "Secrets: Key Vault is canonical. No secret ever appears in Bicep outputs, container images, or deployment history. Managed-identity reads only.",
            "Identity: federated OIDC from GitHub Actions to Azure SP — zero long-lived client secrets stored in CI.",
            "Tenant isolation: SQL-layer enforcement on every read/mutate path. 9 audit gaps closed; parametrized fuzz tests over all three backends in CI.",
            "Cost protection: per-tenant monthly cost ceiling enforced at executor entry — runaway cost limited by config, not by humans noticing.",
            "Network: public ingress in v1.0 with bearer-token auth on every endpoint. VNet integration deferred to v1.1 if security review requires.",
            "Rate limiting per API key (token-bucket; default 60 req/min/key) prevents key-credential abuse from a single source.",
            "Audit log: every action persisted with `api_key_id` for forensics — failures stored with full context in `failures` table.",
            "RBAC + SSO deferred to v1.1+. v1.0 is single-tier bearer-token API keys. Acceptable for internal-only v1.0 scope.",
        ],
    )

    # Slide 21 — What's not in v1.0
    add_table_slide(
        prs,
        title="What's NOT in v1.0",
        subtitle="Explicit deferrals — none are accidental",
        headers=["Capability", "Reason for deferral", "Target"],
        rows=[
            ["RBAC + Azure AD SSO", "Multi-user auth deferred; v1.0 is bearer-token API keys (internal scope)", "v1.1"],
            ["LangGraph TOOL / FUNCTION / SUB_WORKFLOW nodes", "`can_compile` rejects with operator-facing pointer; each PR flips one rejection", "v1.1"],
            ["HTTP streaming for `POST /run?wait=true`", "Polling is fine for batch / dev-team flows; reconsider for interactive UIs", "v1.1"],
            ["DeepEval / Ragas / TruLens integrations", "Defer until a RAG agent ships to production needing those metrics", "v1.1+"],
            ["Multi-region failover", "Customers are single-region; bring active/passive when that changes", "post-v1.x"],
            ["Custom domains + TLS", "ACA's `*.azurecontainerapps.io` URL is fine for v1.0; cert provisioning deferred", "v1.1"],
        ],
    )

    # ──────────────────────────────────────────────────────────────
    # Slide 22 — Section: ROADMAP
    add_section(prs, "Part 4", "Roadmap")

    # Slide 23 — Roadmap: Next 1 week
    add_content_slide(
        prs,
        title="Roadmap — Next 1 week (Tier B polish)",
        subtitle="Highest-leverage investments. Each ~1 day.",
        bullets=[
            "**More agent templates** (extractor, RAG, function-caller) — most customer-visible. Answers \"what can I build?\" with concrete starts. `movate init --template extractor my-thing`.",
            "**`movate logs <run-id> --tail`** — Rich-rendered timeline of stored events. Pairs with Telegram alerts (alert lands → check what happened).",
            "**Privacy redaction (`tracer.redact_io: true`)** — config flag to mask prompt / output in spans for PII-sensitive tenants. Gates real customer onboarding.",
            "**Rubric library** — 5 standard judges (relevance, correctness, faithfulness, safety, tone). Importable by name from `evals/judge.yaml`. Productizes the eval flow.",
        ],
        footer="See: BACKLOG.md Tier B for the full list with effort estimates",
    )

    # Slide 24 — Roadmap: Next 1 month + quarter
    add_two_column_slide(
        prs,
        title="Roadmap — Next 1 month + 1 quarter",
        left_heading="Next 1 month (real features)",
        left_bullets=[
            "HTTP streaming for `/run?wait=true` (SSE) — interactive UI use case",
            "`/run` idempotency by `request_id` — CI retry safety",
            "`workflow_runs` linking table — parent→child run lineage",
            "LangGraph TOOL node compilation — first of the deferred node types",
            "Workflow replay CLI — `movate run --replay <workflow-run-id>`",
        ],
        right_heading="Next 1 quarter (v1.1+)",
        right_bullets=[
            "RBAC + Azure AD SSO — multi-user auth",
            "DeepEval / Ragas integration — RAG-grounding metrics",
            "Multi-region failover — active/passive across regions",
            "Custom domain + TLS — for any customer-facing surface",
            "VNet integration — if security review requires",
        ],
    )

    # Slide 25 — Asks + Q&A
    add_content_slide(
        prs,
        title="Asks + Q&A",
        subtitle="What I need to keep momentum",
        bullets=[
            "**Operator-side: A2P 10DLC brand registration** (~2-3 weeks with The Campaign Registry). Only matters if customer-facing SMS is a real product surface. Telegram covers internal/personal alerts already.",
            "**Azure subscription access for staging + prod** — currently validated on a personal sub. Production deploys need Movate's Azure tenancy.",
            "**Movate corp Azure permissions** for redoing platform validation in the corporate tenant: Application Developer (Entra ID) + Owner on a dedicated `movate-<env>-rg`.",
            "**Customer engagement decision** — which Tier B template is highest priority? Drives which template I scaffold first (extractor / RAG / function-caller).",
            "**Budget signoff** for production deploy — ~$300-500/mo per environment when running.",
            "Open questions / Q&A.",
        ],
        footer="docs/v1.0-overview.md · docs/dev-loop.md · BACKLOG.md · docs/azure-bootstrap.md",
    )

    return prs


if __name__ == "__main__":
    prs = build()
    out = Path("docs/movate-cli-exec-deck.pptx")
    prs.save(out)
    print(f"✓ wrote {out} ({len(prs.slides)} slides)")
