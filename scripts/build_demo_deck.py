"""Build the demo walkthrough deck: "Build a Movate FAQ Agent in 30 min."

Demonstrates the full movate-cli capability surface against a realistic
production pattern: a manager agent that routes questions to role
agents, each with its own knowledge base.

Architecture:
  manager (classifier) ─┬─ "services" → services_expert (Movate svc KB)
                        └─ "cli"      → cli_expert (movate-cli KB)

This exercises: typed agents, multi-step workflows, conditional edges
(via runtime: langgraph), eval gates, multi-model bench, deploy, and
production submission with Telegram alerts.

Audience: engineers / SAs / PLs who have seen the onboarding deck and
want a worked example before building their own agent.

Usage:
    uv run python scripts/build_demo_deck.py

Output: docs/movate-cli-demo-deck.pptx
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt

# Same palette as the other two decks so they read as a set
INK = RGBColor(0x1A, 0x1A, 0x1A)
INK_DIM = RGBColor(0x55, 0x55, 0x55)
ACCENT = RGBColor(0x00, 0x6E, 0xB8)
ACCENT_LITE = RGBColor(0x00, 0xA8, 0xE8)
GREEN = RGBColor(0x2E, 0x86, 0x36)
AMBER = RGBColor(0xC1, 0x7B, 0x00)
RED = RGBColor(0xB0, 0x2A, 0x2A)
CODE_BG = RGBColor(0x1E, 0x1E, 0x2E)
CODE_FG = RGBColor(0xE6, 0xE6, 0xE6)
BG = RGBColor(0xFF, 0xFF, 0xFF)
BG_PANEL = RGBColor(0xF6, 0xF7, 0xF9)


# -----------------------------------------------------------------------------
# Slide builders
# -----------------------------------------------------------------------------


def add_title_slide(prs, title, subtitle, kicker):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_text(slide, kicker.upper(), left=0.6, top=2.2, width=12.0, height=0.5,
              size=14, bold=True, color=ACCENT)
    _add_text(slide, title, left=0.6, top=2.7, width=12.0, height=1.6,
              size=52, bold=True, color=INK)
    _add_text(slide, subtitle, left=0.6, top=4.3, width=12.0, height=1.5,
              size=22, color=INK_DIM)
    _add_text(slide, "Hands-on demo · pairs with the onboarding deck",
              left=0.6, top=6.6, width=12.0, height=0.5,
              size=14, color=INK_DIM)


def add_section(prs, label, title):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_filled_rect(slide, left=0.0, top=0.0, width=13.33, height=7.5, color=BG_PANEL)
    _add_text(slide, label.upper(), left=0.6, top=2.8, width=12.0, height=0.5,
              size=14, bold=True, color=ACCENT)
    _add_text(slide, title, left=0.6, top=3.3, width=12.0, height=1.5,
              size=42, bold=True, color=INK)


def add_content_slide(prs, title, bullets, *, subtitle=None, footer=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_filled_rect(slide, left=0.6, top=1.05, width=12.13, height=0.04, color=ACCENT)
    _add_text(slide, title, left=0.6, top=0.4, width=12.13, height=0.7,
              size=28, bold=True, color=INK)
    body_top = 1.3
    if subtitle:
        _add_text(slide, subtitle, left=0.6, top=1.2, width=12.13, height=0.4,
                  size=14, color=INK_DIM)
        body_top = 1.65
    _add_bullets(slide, bullets, left=0.6, top=body_top, width=12.13, height=5.5,
                 size=17, color=INK)
    if footer:
        _add_text(slide, footer, left=0.6, top=6.95, width=12.13, height=0.4,
                  size=11, color=INK_DIM)


def add_code_slide(prs, title, *, subtitle=None, code, explanation=None, footer=None,
                   code_lang="bash"):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_filled_rect(slide, left=0.6, top=1.05, width=12.13, height=0.04, color=ACCENT)
    _add_text(slide, title, left=0.6, top=0.4, width=12.13, height=0.7,
              size=28, bold=True, color=INK)
    body_top = 1.3
    if subtitle:
        _add_text(slide, subtitle, left=0.6, top=1.2, width=12.13, height=0.4,
                  size=14, color=INK_DIM)
        body_top = 1.65

    code_lines = code.strip().split("\n")
    code_h = max(1.5, min(4.4, 0.28 * len(code_lines) + 0.4))
    _add_filled_rect(slide, left=0.6, top=body_top, width=12.13, height=code_h,
                     color=CODE_BG)
    _add_code_text(slide, code.strip(), left=0.8, top=body_top + 0.15,
                   width=11.7, height=code_h - 0.3)

    if explanation:
        _add_bullets(slide, explanation,
                     left=0.6, top=body_top + code_h + 0.2, width=12.13,
                     height=6.9 - (body_top + code_h + 0.2),
                     size=14, color=INK)

    if footer:
        _add_text(slide, footer, left=0.6, top=6.95, width=12.13, height=0.4,
                  size=11, color=INK_DIM)


def add_mermaid_text_slide(prs, title, *, subtitle=None, ascii_diagram, caption=None):
    """ASCII-art architecture diagram. PowerPoint renders monospace cleanly;
    no need for actual Mermaid rendering."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_filled_rect(slide, left=0.6, top=1.05, width=12.13, height=0.04, color=ACCENT)
    _add_text(slide, title, left=0.6, top=0.4, width=12.13, height=0.7,
              size=28, bold=True, color=INK)
    body_top = 1.3
    if subtitle:
        _add_text(slide, subtitle, left=0.6, top=1.2, width=12.13, height=0.4,
                  size=14, color=INK_DIM)
        body_top = 1.65

    diagram_h = 4.5
    _add_filled_rect(slide, left=0.6, top=body_top, width=12.13, height=diagram_h,
                     color=BG_PANEL)
    _add_code_text(slide, ascii_diagram, left=0.9, top=body_top + 0.2,
                   width=11.5, height=diagram_h - 0.4, color=INK)

    if caption:
        _add_text(slide, caption,
                  left=0.6, top=body_top + diagram_h + 0.2,
                  width=12.13, height=0.8,
                  size=14, color=INK_DIM)


# -----------------------------------------------------------------------------
# Primitives
# -----------------------------------------------------------------------------


def _add_text(slide, text, *, left, top, width, height, size=16, bold=False,
              italic=False, color=INK, align="left"):
    tx = slide.shapes.add_textbox(Inches(left), Inches(top),
                                   Inches(width), Inches(height))
    tf = tx.text_frame
    tf.word_wrap = True
    tf.text = text
    p = tf.paragraphs[0]
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.italic = italic
    p.font.color.rgb = color
    if align == "center":
        from pptx.enum.text import PP_ALIGN  # noqa: PLC0415
        p.alignment = PP_ALIGN.CENTER


def _add_bullets(slide, items, *, left, top, width, height, size, color):
    tx = slide.shapes.add_textbox(Inches(left), Inches(top),
                                   Inches(width), Inches(height))
    tf = tx.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.space_after = Pt(6)


def _add_code_text(slide, code, *, left, top, width, height, color=CODE_FG):
    tx = slide.shapes.add_textbox(Inches(left), Inches(top),
                                   Inches(width), Inches(height))
    tf = tx.text_frame
    tf.word_wrap = False
    tf.text = code
    for p in tf.paragraphs:
        p.font.name = "Menlo"
        p.font.size = Pt(12)
        p.font.color.rgb = color
        p.space_after = Pt(0)


def _add_filled_rect(slide, *, left, top, width, height, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                   Inches(left), Inches(top),
                                   Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    shape.shadow.inherit = False


# -----------------------------------------------------------------------------
# Deck content
# -----------------------------------------------------------------------------


def build():
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    # ──────────────────────────────────────────────────────────────
    # 1 — Title
    add_title_slide(
        prs,
        kicker="Demo walkthrough",
        title="Build a Movate FAQ agent",
        subtitle="Manager + 2 role agents · routing, eval, bench, deploy in 30 min",
    )

    # 2 — What you'll build
    add_content_slide(
        prs,
        title="What you'll build",
        subtitle="A real-world pattern: a manager agent that routes questions to specialists",
        bullets=[
            "**Manager agent** — classifies an incoming question as \"services\" or \"cli\".",
            "**Services-expert agent** — answers questions about Movate's service offerings. KB is a list of Movate services baked into the prompt.",
            "**CLI-expert agent** — answers questions about movate-cli itself. KB is excerpts from the docs/ folder.",
            "**Workflow** — a single `workflow.yaml` wires the three agents together. Conditional edges route based on the manager's classification.",
            "**Eval dataset** — 30 test questions (mix of services + CLI topics). Gate on 90%+ routing accuracy.",
            "**Bench** — compare 3 model choices for the manager role. Cheapest model that still routes accurately wins.",
            "**Deploy** — push to Movate's Azure runtime. Submit jobs from anywhere. Telegram pings your phone when answers land.",
        ],
        footer="By the end: you've exercised typed agents, conditional workflows, eval gates, multi-model bench, deploy, and prod submission — the full v1.0 surface.",
    )

    # 3 — Architecture diagram
    add_mermaid_text_slide(
        prs,
        title="The architecture",
        subtitle="One workflow, three agents, conditional routing",
        ascii_diagram="""
                    ┌──────────────┐
                    │   USER       │
                    │  question    │
                    └──────┬───────┘
                           │
                           ▼
                ┌──────────────────────┐
                │    manager agent     │  (classifies the question)
                │  ──────────────────  │
                │  outputs:            │
                │    classification    │
                │      ∈ {services,    │
                │         cli}         │
                └──────────┬───────────┘
                           │
                  ┌────────┴─────────┐
                  │   conditional    │
                  │      edge        │
                  └────────┬─────────┘
                           │
              ┌────────────┴────────────┐
              │                         │
              ▼ classification=services ▼ classification=cli
   ┌────────────────────┐     ┌────────────────────┐
   │ services_expert    │     │   cli_expert       │
   │ KB: Movate services│     │ KB: movate-cli docs│
   └─────────┬──────────┘     └──────────┬─────────┘
             │                            │
             └─────────────┬──────────────┘
                           ▼
                  ┌──────────────────┐
                  │   final answer   │
                  └──────────────────┘
""",
        caption="Three typed agents wired by one workflow.yaml. Conditional edges (movate v1.1 / LangGraph runtime) drive the routing.",
    )

    # 4 — Why this is a good first demo
    add_content_slide(
        prs,
        title="Why this is a great first demo to build",
        subtitle="Exercises every v1.0 capability against a realistic production pattern",
        bullets=[
            "**Routing / triage is the most common multi-agent pattern.** Customer support, ticket dispatch, internal Q&A bots — they all start with \"figure out who should answer this.\"",
            "**Easy to extend.** Add a third role agent for billing? One YAML edit + one new prompt. The manager doesn't change.",
            "**Tests the eval surface meaningfully.** Routing accuracy is binary per case (right or wrong classification) — perfect for `--gate 0.9` regression detection.",
            "**Bench has a clear winner condition.** Manager is small + frequent: you want the cheapest model that still routes accurately. `movate bench` quantifies the trade-off.",
            "**Shows off conditional edges** — the v1.1 LangGraph feature that distinguishes movate workflows from simple linear chains.",
            "**Real production shape.** Every customer engagement at Movate that ships agents will look something like this — multiple specialized agents, a router, a knowledge base per specialist.",
        ],
    )

    # 5 — Prereqs
    add_code_slide(
        prs,
        title="Before you start",
        subtitle="What you need on your machine",
        code="""# movate-cli installed
movate --version    # expect: 1.0.0 or later

# At least 2 provider API keys in .env (or env vars)
cat .env
# OPENAI_API_KEY=sk-proj-...
# ANTHROPIC_API_KEY=sk-ant-api03-...

# Optional but recommended: register Movate's prod deployment target
# (skip if running purely local for now)
movate config list-targets""",
        explanation=[
            "You'll use OpenAI for the manager (cheap classifier role) and Anthropic for the role agents (longer-form answers).",
            "Movate's runtime target is optional — you can build + eval + bench entirely local first, deploy only when you're happy.",
            "Estimated provider spend for this whole walkthrough (3 dev runs + eval + bench): ~$0.15.",
        ],
        footer="If movate isn't installed yet, see the onboarding deck (slides 5-8) or run: uv tool install --from git+https://github.com/Movate/movate-cli movate-cli",
    )

    # ──────────────────────────────────────────────────────────────
    # Reference — Anatomy of a movate agent
    #
    # Five slides explaining what each file in the canonical layout
    # does, before the audience hits a real agent.yaml in Part 1.
    # Pedagogical, not technical: WHY each piece exists, WHAT it does,
    # one concrete example. New presenters can lean on these slides as
    # the "schema 101" segment.
    # ──────────────────────────────────────────────────────────────
    add_section(prs, "Reference", "Anatomy of a movate agent")

    # Agent layout overview — annotated tree
    add_code_slide(
        prs,
        title="The canonical agent layout",
        subtitle="Four files. Each one does one job. Movate enforces the contract between them.",
        code="""agents/my-agent/
├── agent.yaml          ← the contract: name, model, schema refs, budget
├── prompt.md           ← the instruction template (with {{input.x}} vars)
├── schema/
│   ├── input.json      ← JSON Schema — what callers must send
│   └── output.json     ← JSON Schema — what the model must return
└── evals/
    ├── dataset.jsonl   ← test cases (one per line)
    └── judge.yaml      ← how to score: exact match or LLM-as-judge""",
        explanation=[
            "**Separation of concerns** — the YAML is a contract, the prompt is plain English, the schemas are types, the evals are the safety net. Each piece is independently version-controlled and review-able.",
            "**`movate init my-agent` scaffolds all of this** — you don't write it from a blank page. Templates: `faq`, `classifier`, `extractor`, `chatbot`.",
            "**`movate validate` checks all five together** — schema parses? prompt linter clean? policy compliant? cost forecast under cap?",
        ],
        footer="The next four slides explain each file in detail.",
    )

    # agent.yaml — the contract
    add_code_slide(
        prs,
        title="agent.yaml — the contract",
        subtitle="Identity, model choice, schema references, budget cap. Validated by Pydantic.",
        code="""api_version: movate/v1     # YAML schema version
name: faq-agent             # used in run records, logs, deploys
version: 0.1.0              # semver — bump triggers re-eval gate

model:
  provider: openai/gpt-4o-mini-2024-07-18    # LiteLLM-style
  params:
    temperature: 0.2
    max_tokens: 1024
  fallback:                                    # optional retry chain
    - provider: anthropic/claude-haiku-4-5-20251001

prompt: ./prompt.md                            # path to instructions
schema:
  input: ./schema/input.json                   # JSON Schema for input
  output: ./schema/output.json                 # JSON Schema for output

budget:
  max_cost_usd_per_run: 0.05    # safety net — hard cap per call""",
        explanation=[
            "**Provider string is LiteLLM format** (`<vendor>/<model>`). Switching from OpenAI to Anthropic = one line. No code change.",
            "**Fallback chain** — if the primary errors (rate limit, timeout, content filter), executor automatically retries on the fallback. Both rows in your run record.",
            "**Budget cap is per-run, not per-month** — protects against pathological inputs that produce 50k-token answers. Combined with the project-level cost ceiling in `movate.yaml`.",
        ],
    )

    # prompt.md — the instruction template
    add_code_slide(
        prs,
        title="prompt.md — the instruction template",
        subtitle="Plain Markdown with Jinja2 placeholders. Diff-reviewable, version-controlled.",
        code="""You are a helpful FAQ assistant for Movate.

Answer the user's question using the knowledge base below. If the
answer isn't in the KB, say so honestly — don't make up facts.

## Knowledge base
- Movate is a digital engineering services company.
- We provide consulting, managed services, and AI/ML solutions.
- Headquartered in Plano, TX. Founded 2002.

## User question
{{ input.question }}

## Response format
Respond with JSON matching this schema:
  { "answer": "<your answer>", "confident": true | false }""",
        explanation=[
            "**Jinja2 substitution** — `{{ input.<field> }}` pulls values from the validated request input. Refs to fields not in the input schema fail at `movate validate`.",
            "**Linter rules** — empty prompt, missing JSON instruction, undeclared input refs, tiny prompt. Runs on every save (`movate watch`) and in CI.",
            "**Prompt hash is persisted in every RunRecord** — proves exactly which prompt version produced each result. Critical for incident replay.",
        ],
    )

    # schema/ folder — input + output JSON Schemas
    add_code_slide(
        prs,
        title="schema/ — input + output contracts",
        subtitle="JSON Schema (draft 2020-12). Pydantic validates at request + response time.",
        code="""# schema/input.json — what callers must send
{
  "type": "object",
  "properties": {
    "question": { "type": "string", "minLength": 1 }
  },
  "required": ["question"],
  "additionalProperties": false
}

# schema/output.json — what the model must return
{
  "type": "object",
  "properties": {
    "answer":    { "type": "string" },
    "confident": { "type": "boolean" }
  },
  "required": ["answer", "confident"],
  "additionalProperties": false
}""",
        explanation=[
            "**Input schema rejects bad calls before any LLM spend** — missing `question`? `additionalProperties: false`? 422 at the door, $0.00 cost.",
            "**Output schema is the safety net** — model returned `{ \"answr\": ... }` (typo)? Executor catches it, retries (per `model.fallback`), and if retries exhausted records a typed `SchemaError` failure.",
            "**Enums in output = guarantees**. The classifier agent uses `\"classification\": { \"enum\": [\"services\", \"cli\"] }` — the model literally cannot return a third option without a hard validation fail.",
        ],
    )

    # evals/ folder — dataset + judge
    add_code_slide(
        prs,
        title="evals/ — the quality net",
        subtitle="Test cases + scoring method. Gateable in CI; baseline-diffable across versions.",
        code="""# evals/dataset.jsonl — one test case per line
{"input": {"question": "What does Movate do?"}, "expected": {"confident": true}}
{"input": {"question": "Where is HQ?"}, "expected": {"confident": true}}
{"input": {"question": "What is quantum chromodynamics?"}, "expected": {"confident": false}}
# ...30 cases total

# evals/judge.yaml — how each case is scored
method: llm_judge                  # or `exact_match` for deterministic outputs
model:
  provider: anthropic/claude-haiku-4-5-20251001   # judge != tested family
rubric: |
  Score the answer 0.0-1.0. Full credit if the answer is factually
  accurate AND the `confident` field correctly reflects whether the
  KB had the answer. Half credit if accurate but mis-confident.""",
        explanation=[
            "**Exact match** for classifiers + extractors (deterministic outputs). **LLM-as-judge** for free-form prose. Cross-family enforcement: judge model ≠ tested model's vendor (prevents same-family bias).",
            "**`movate eval --gate 0.9`** → exit 1 if pass rate < 90%. CI-gateable. `--gate-mode mean|min|p10` for different risk profiles.",
            "**`movate eval --baseline <id>`** → diffs scores against a stored eval. Drift detection: regression past `--regression-tolerance` fails CI. Same pattern works for `movate bench --baseline`.",
        ],
    )

    # ──────────────────────────────────────────────────────────────
    # 6 — Section: Build the agents
    add_section(prs, "Part 1", "Build the three agents")

    # 7 — Project layout
    add_code_slide(
        prs,
        title="Step 1 — Project layout",
        subtitle="One workflow project containing three agent subdirs",
        code="""mkdir movate-faq && cd movate-faq

# Scaffold each agent in its own subdir
movate init agents/manager -t classifier
movate init agents/services -t faq
movate init agents/cli -t faq

# Files we'll create by hand:
#   workflow.yaml   — wires the three agents
#   state.json      — typed state schema
#   evals/dataset.jsonl — test cases

# Final layout
find . -type d | sort""",
        explanation=[
            "**`classifier` template** for the manager — it ships with an exact-match judge config since classifications are deterministic (services or cli, not both).",
            "**`faq` template** for the role agents — ships with an LLM-as-judge config since their answers are free-form prose.",
            "`workflow.yaml` is where you'll wire them together. The next slides build it up.",
        ],
    )

    # 8 — Manager: scaffold + prompt
    add_code_slide(
        prs,
        title="Step 2 — The manager agent",
        subtitle="A cheap, fast classifier. No knowledge base; just routing logic.",
        code="""# agents/manager/agent.yaml
name: manager
version: 0.1.0
model:
  provider: openai/gpt-4o-mini-2024-07-18    # cheap; classification is easy
  params: { temperature: 0 }                  # deterministic
prompt: ./prompt.md
schema:
  input: ./schema/input.json     # { question: string }
  output: ./schema/output.json   # { classification: 'services' | 'cli' }
budget:
  max_cost_usd_per_run: 0.001    # hard cap

# agents/manager/prompt.md (excerpt)
Classify the user's question. Output JSON with a "classification"
field. Use "services" for questions about Movate's offerings
(consulting, managed services, training). Use "cli" for questions
about the movate-cli toolkit (commands, features, deploys, evals).

Question: {{ input.question }}""",
        explanation=[
            "**Temperature 0** for the classifier — same question should always classify the same way. Variance here is bad.",
            "**Output schema with an enum** — the JSON Schema validator rejects anything other than \"services\" or \"cli\". Manager can't accidentally output a third category.",
            "**$0.001 budget cap** — manager is the high-frequency hot path. A misbehaved manager that runs a $0.01 call per question would be unfair to downstream cost.",
        ],
    )

    # 9 — Services expert: scaffold + KB
    add_code_slide(
        prs,
        title="Step 3 — The services expert",
        subtitle="Answers Movate-service questions; KB is embedded in the prompt",
        code="""# agents/services/agent.yaml
name: services
model:
  provider: anthropic/claude-haiku-4-5-20251001   # great for prose
  params: { temperature: 0.3 }                     # slight variance OK
prompt: ./prompt.md
schema:
  input: ./schema/input.json    # { question: string }
  output: ./schema/output.json  # { answer: string, sources: [string] }
budget:
  max_cost_usd_per_run: 0.05

# agents/services/prompt.md (excerpt — KB is inline for demo)
You are a Movate services expert. Use the knowledge base below to
answer questions about Movate's offerings.

# Knowledge base
- **Digital Engineering**: cloud transformation, app modernization, ...
- **Managed Services**: 24/7 infrastructure ops, security ops center, ...
- **Industry Solutions**: retail, banking, healthcare verticals, ...
- (...etc)

Question: {{ input.question }}
Answer with citations from the KB above. Output {answer, sources}.""",
        explanation=[
            "**KB inline in prompt** — fine for a demo with a small fixed KB. Production would swap to a `kb_search` tool (movate's `@tool` decorator) backed by pgvector or Azure AI Search.",
            "**Anthropic Claude** for prose — handles longer-form answers more reliably than the GPT-4o-mini you used for the manager. Per-role model selection is normal.",
            "**Output schema requires `sources`** — forces the model to ground its answer in the KB. If it tries to invent a source not in the KB, eval-as-judge will catch it.",
        ],
    )

    # 10 — CLI expert: scaffold + KB
    add_code_slide(
        prs,
        title="Step 4 — The CLI expert",
        subtitle="Same pattern, different KB",
        code="""# agents/cli/prompt.md (excerpt)
You are a movate-cli expert. Answer questions about how to use the
toolkit: scaffolding agents, running evals, deploying, etc.

# Knowledge base
- `movate init` — scaffolds a new agent from a template (faq,
  summarizer, classifier, or generic agent_init).
- `movate run` — single-shot execution. --mock for no-API-call.
- `movate eval` — dataset eval with --gate score regression.
- `movate bench` — multi-model comparison.
- `movate deploy` — push to Azure Container Apps.
- (...etc — populate from docs/dev-loop.md)

# Examples of high-level questions
- "How do I create a new agent?"
- "What is a workflow?"
- "How do I deploy to production?"

# Examples of technical questions
- "How does --baseline detect drift?"
- "What's the difference between eval and bench?"
- "How does conditional routing work in workflow.yaml?"

Question: {{ input.question }}""",
        explanation=[
            "**KB sourced from the repo's docs/** — `docs/dev-loop.md`, README, command help text. For a richer KB you'd extract per-section markdown and embed; that's a follow-up.",
            "**\"Examples of...\" sections in the prompt** prime the model to handle both high-level (\"what's a workflow?\") and technical (\"how does --baseline diff?\") questions. Few-shot in disguise.",
            "**Same output schema** as services expert — keeps the downstream code one branch instead of two.",
        ],
    )

    # ──────────────────────────────────────────────────────────────
    # 11 — Section: Wire the workflow
    add_section(prs, "Part 2", "Wire the workflow")

    # 12 — workflow.yaml with conditional edges
    add_code_slide(
        prs,
        title="Step 5 — The workflow.yaml",
        subtitle="One file describes the entire routing topology",
        code="""# movate-faq/workflow.yaml
api_version: movate/v1
kind: Workflow
name: movate-faq
version: 0.1.0
description: Routes questions to services or CLI experts.

# Use LangGraph runtime for conditional edges + checkpointing
runtime: langgraph
checkpointer: memory    # use sqlite/postgres in production

state_schema: ./state.json
entrypoint: manager

nodes:
  - { id: manager,          type: agent, ref: ./agents/manager  }
  - { id: services_expert,  type: agent, ref: ./agents/services }
  - { id: cli_expert,       type: agent, ref: ./agents/cli      }

edges:
  - from: manager
    to: services_expert
    kind: conditional
    when: "$.classification == 'services'"
  - from: manager
    to: cli_expert
    kind: conditional
    when: "$.classification == 'cli'"
  - from: manager
    to: cli_expert
    kind: conditional
    when: null        # default branch — required + must be last""",
        explanation=[
            "**`runtime: langgraph`** unlocks conditional edges. The default homegrown runner only handles linear DAGs.",
            "**JSONPath-like DSL** in the `when:` clauses — supports `==`, `!=`, `<`, `>`, `&&`, `||`, `in [...]`. No `eval()`, no third-party dep.",
            "**Default branch (`when: null`) is required and must be last** — covers the \"manager misclassified\" or \"unknown intent\" case. We route to cli_expert as a safe fallback.",
        ],
    )

    # 13 — state.json
    add_code_slide(
        prs,
        title="Step 6 — The state schema",
        subtitle="Typed state flows through the workflow; each node adds to it",
        code="""# movate-faq/state.json
{
  "$schema": "https://json-schema.org/draft/2020-12/schema",
  "type": "object",
  "required": ["question"],
  "properties": {
    "question":       { "type": "string" },
    "classification": { "type": "string", "enum": ["services", "cli"] },
    "answer":         { "type": "string" },
    "sources":        { "type": "array", "items": { "type": "string" } }
  }
}""",
        explanation=[
            "**`question` is required** as the initial state. Any user submission must provide it.",
            "**Manager fills `classification`** — its output is merged into state, then the conditional edge reads `$.classification` to route.",
            "**Role agents fill `answer` + `sources`** — final workflow state has all four fields.",
            "**Enum constrains `classification`** — schema validates the manager's output before the conditional edge fires.",
        ],
    )

    # 14 — First local run + --node-trace
    add_code_slide(
        prs,
        title="Step 7 — First local run",
        subtitle="`movate run` walks the workflow end-to-end · `--node-trace` shows state evolution",
        code="""# Run with a CLI question — expect routing → cli_expert
movate run ./movate-faq \\
    '{"question": "How do I deploy an agent to Azure?"}' \\
    --node-trace

# Output (abbreviated):
#  ✓ workflow run mvt-9ad12fe... — success (1.2s, $0.018)
#
#  Node-by-node state trace
#  initial state: {"question": "How do I deploy an agent to Azure?"}
#
#  1. manager
#     output:    {"classification": "cli"}
#     + added:   classification = "cli"
#
#  2. cli_expert
#     output:    {"answer": "Use `movate deploy --target prod`...",
#                 "sources": ["docs/dev-loop.md"]}
#     + added:   answer = "..."
#                sources = [...]""",
        explanation=[
            "**`--node-trace`** is the workflow-debugging killer feature. Shows the state delta added by each node — instantly spot which node corrupted a key.",
            "**Run cost surfaced inline** ($0.018 here) — manager + role agent. The manager was $0.0003, role agent was $0.018. Manager's cheap; that's by design.",
            "**Try a services question next**: `'{\"question\": \"What does Movate's managed services offer?\"}'` — you'll see the workflow route to services_expert.",
        ],
    )

    # ──────────────────────────────────────────────────────────────
    # 15 — Section: Quality + cost
    add_section(prs, "Part 3", "Quality + cost gates")

    # 16 — Build eval dataset
    add_code_slide(
        prs,
        title="Step 8 — Build an eval dataset",
        subtitle="Mix of services and CLI questions; expected classification per case",
        code="""# movate-faq/evals/dataset.jsonl
{"input": {"question": "What is Movate's digital engineering practice?"},
 "expected": {"classification": "services"}}
{"input": {"question": "How do I run movate eval --gate 0.7?"},
 "expected": {"classification": "cli"}}
{"input": {"question": "Tell me about Movate's industry solutions"},
 "expected": {"classification": "services"}}
{"input": {"question": "What's the difference between --replay and --baseline?"},
 "expected": {"classification": "cli"}}
# ... 30 cases total, 15 of each type

# movate-faq/evals/judge.yaml
method: exact_match     # routing is binary — no LLM judge needed
field: classification   # only score on the manager's output""",
        explanation=[
            "**Mix the two question types ~50/50** — catches a manager that always answers \"cli\" with 50% accuracy. Skewed datasets miss skewed biases.",
            "**Exact-match judge** for classification — no LLM-as-judge needed; the manager's output is one of two strings.",
            "**Eval on just `classification` field** — we're testing the manager's routing, not the role agents' prose. Separate concerns get separate tests.",
        ],
    )

    # 17 — Run eval, gate on accuracy
    add_code_slide(
        prs,
        title="Step 9 — Run eval, gate on routing accuracy",
        subtitle="`--gate 0.9` fails CI on routing regression",
        code="""# Eval the workflow with the manager's routing as the gated output
movate eval ./movate-faq --gate 0.9 --runs 3

# Sample output (Rich table):
#  ┌────────┬───────────────┬─────────────────┬────────────┐
#  │ case # │ expected      │ got             │ score      │
#  ├────────┼───────────────┼─────────────────┼────────────┤
#  │ 1      │ services      │ services        │ ✓ 1.000    │
#  │ 2      │ cli           │ cli             │ ✓ 1.000    │
#  │ 3      │ services      │ cli             │ ✗ 0.000    │  ← miss
#  │ ...    │               │                 │            │
#  └────────┴───────────────┴─────────────────┴────────────┘
#
#  mean_score: 0.933 (28/30 correct)  ✓ above gate 0.9
#
#  Saved as eval_id=ev-abc-123... (use --baseline ev-abc-123 to detect drift)""",
        explanation=[
            "**3 runs/case** smooths the rare \"GPT-4o-mini hiccup\" non-determinism. Gate-mode `mean` (the default) averages.",
            "**93.3% mean accuracy ≥ 0.9 gate** — exits 0; CI passes. A regression to 85% would exit 1 and block the merge.",
            "**Stored `eval_id`** — pass it to the next eval as `--baseline ev-abc-123` to catch \"this used to work\" regressions specifically.",
        ],
    )

    # 18 — Bench the manager
    add_code_slide(
        prs,
        title="Step 10 — Bench the manager's model choice",
        subtitle="Which provider gives best routing accuracy for the cost?",
        code="""# Bench just the manager agent across 3 models
movate bench ./agents/manager \\
    '{"question": "How do I deploy an agent?"}' \\
    --runs 5 \\
    -m openai/gpt-4o-mini-2024-07-18 \\
    -m anthropic/claude-haiku-4-5-20251001 \\
    -m google/gemini-2-flash

# Sample output (Rich table):
#  ┌──────────────────────────────────┬──────────┬─────┬─────┬───────┐
#  │ model                             │ cost/run │ p50 │ p95 │ score │
#  ├──────────────────────────────────┼──────────┼─────┼─────┼───────┤
#  │ openai/gpt-4o-mini-2024-07-18    │ $0.0003  │ 280 │ 410 │ 1.000 │
#  │ anthropic/claude-haiku-4-5-...    │ $0.0008  │ 510 │ 690 │ 1.000 │
#  │ google/gemini-2-flash             │ $0.0001  │ 340 │ 450 │ 1.000 │
#  └──────────────────────────────────┴──────────┴─────┴─────┴───────┘
#
#  Saved as bench_id=bn-xyz-789...""",
        explanation=[
            "**All three at 100% accuracy** on a single test question — manager's task is easy. The cost tie-breaker picks Gemini-2-flash ($0.0001/run).",
            "**Run the same bench with a full 30-case dataset** for production-grade signal: `movate bench ./agents/manager --dataset evals/dataset.jsonl`.",
            "**Save the bench_id**. Next time you tweak the manager prompt, run with `--baseline bn-xyz-789 --regression-tolerance 0.05` — CI fails if a model regresses past 5%.",
        ],
    )

    # ──────────────────────────────────────────────────────────────
    # 19 — Section: Ship it
    add_section(prs, "Part 4", "Ship it to production")

    # 20 — Deploy
    add_code_slide(
        prs,
        title="Step 11 — Deploy to Azure",
        subtitle="Workflow + 3 agents + KB go up; ACA scales worker on queue depth",
        code="""# Make sure your target is registered (one-time)
movate config show

# Build container + push to ACR + roll new revision on both apps
movate deploy --target prod

# Output (abbreviated):
#  ✓ Building image movate:0.5.0-9e98a1e via az acr build... done (54s)
#  ✓ Updating movate-prod-api... revision movate-prod-api--rev-xyz live
#  ✓ Updating movate-prod-worker... revision movate-prod-worker--rev-xyz live
#  ✓ /healthz reports v0.5.0-9e98a1e on both apps
#
#  Deploy successful in 87s. Run movate doctor --target prod for sanity.

# Sanity-check the deploy
movate doctor --target prod
# Expect: 9 green checks (az login, sub, RG, ACR, both apps, /healthz)""",
        explanation=[
            "**`az acr build`** does the image build server-side — no Docker needed on your laptop. ~1 min for an incremental change; ~2 min cold.",
            "**Revision roll is atomic** — ACA holds traffic on the previous revision until the new one's `/healthz` returns the new version string.",
            "**Rollback is one command**: `movate deploy --target prod --skip-build --image-tag movate:0.5.0-<prev-sha>`. No re-build needed.",
        ],
    )

    # 21 — Production submit + Telegram
    add_code_slide(
        prs,
        title="Step 12 — Use it in production",
        subtitle="Submit jobs from anywhere · Telegram pings your phone",
        code="""# Submit a workflow run against the deployed runtime
movate submit movate-faq \\
    '{"question": "What is Movate'\\''s digital engineering practice?"}' \\
    --target prod \\
    --wait

# Output:
#  ✓ workflow_run mvt-prod-456... — success (2.4s, $0.024)
#  classification: services
#  answer: Movate's digital engineering practice covers...
#  sources: ["Digital Engineering", "Industry Solutions"]
#
#  📱 Your phone dings via Telegram bot

# Fire-and-forget mode + poll later
movate submit movate-faq '{"question": "..."}' --target prod
movate jobs show <job-id>

# Recent jobs (paginated)
movate jobs list --target prod --status success --limit 10""",
        explanation=[
            "**`--wait`** blocks until terminal — useful for one-off interactive use. Fire-and-forget by default for scripted submission.",
            "**Telegram alerts** require one-time setup (see onboarding deck slide 19). Once wired, every job your account submits pings your phone.",
            "**Multi-tenant safe** — your `tenant_id` is encoded in your API key. You can never see another tenant's jobs or trigger their workflows.",
        ],
    )

    # ──────────────────────────────────────────────────────────────
    # 22 — Section: Recap + going further
    add_section(prs, "Part 5", "Recap + going further")

    # 23 — What you exercised
    add_content_slide(
        prs,
        title="What you just exercised",
        subtitle="The full v1.0 movate-cli capability surface against one realistic demo",
        bullets=[
            "**3 typed agents** with separate model choices (cheap classifier, prose-strong role agents).",
            "**Conditional workflow** with the LangGraph backend — routing via a JSONPath-like DSL on the manager's output.",
            "**Eval gate** on routing accuracy — `--gate 0.9` fails CI on regression. Per-case results in a Rich table.",
            "**Multi-model bench** with stored `bench_id` for drift detection — cheapest-model-that-routes-correctly is the win condition.",
            "**Cost ceilings + per-tenant budgets** enforced at executor entry. Manager has $0.001/run cap; role agents have $0.05/run.",
            "**Deploy + revision roll** in one command. ACA worker scales on queue depth via KEDA Postgres scaler.",
            "**Production submission** with Telegram alerts on completion. Tenant-isolated; cross-tenant invisible.",
            "**One file (`workflow.yaml`)** describes the whole architecture — version-controlled, PR-reviewable, env-portable.",
        ],
    )

    # 24 — Extensions
    add_content_slide(
        prs,
        title="Going further — extension ideas",
        subtitle="Drop-in changes that scale this from demo to real product",
        bullets=[
            "**Swap inline KB for retrieval-augmented**: replace the prompt-embedded KB with a `@tool`-decorated `kb_search(query)` function backed by pgvector or Azure AI Search. Each role agent calls it per question. Movate-cli's tool registry handles the JSON-schema plumbing.",
            "**Add more role agents**: HR-policy expert, billing expert, contract-lookup. Add a node + a conditional edge — manager's classification enum grows.",
            "**Add HITL approval** before answering sensitive questions: insert a `type: human` node between manager and role agents. External system (Slack approval bot, email, etc.) decides whether to continue.",
            "**Stream answers back via SSE**: opt-in `POST /run?wait=true` with HTTP streaming (Tier B follow-up). Useful for an interactive chat UI on top.",
            "**Multi-language**: have manager classify language first, route to language-specific role agents. Movate's deploy story is single-region today; multi-region is a v1.1+ deferred item.",
            "**Tighter cost gates**: set per-tenant monthly cost ceiling with `movate tenants set-budget --monthly-usd-limit 50`. Caps runaway spend regardless of how many jobs land.",
        ],
        footer="Each extension is additive — no need to rewrite the demo to upgrade incrementally.",
    )

    # 25 — Closing — try it now
    add_content_slide(
        prs,
        title="Try it now",
        subtitle="Start with the smallest possible version and grow it",
        bullets=[
            "**Minimum viable demo (~30 min):** scaffold the three agents, write workflow.yaml, run with one example question. No KB, no eval, no deploy.",
            "**Add KB (~30 min):** populate role agent prompts with a real KB excerpt. Run with 5 different questions; verify routing is correct.",
            "**Add eval (~30 min):** build the 30-case dataset, run `movate eval --gate 0.9`. Iterate the manager prompt until it passes.",
            "**Add bench (~15 min):** compare 3 manager models. Pick the cheapest that still routes accurately.",
            "**Deploy (~10 min):** one command, end-to-end. Hit your phone with a Telegram ping.",
            "Total: about 2 hours from \"empty directory\" to \"production-deployed multi-agent system with quality gates.\"",
            "When you ship one — share in #movate-cli. Other teams' demos are the best documentation of what's possible.",
        ],
        footer="Pair this deck with the onboarding deck (commands reference) and docs/dev-loop.md (deeper walkthrough).",
    )

    return prs


if __name__ == "__main__":
    prs = build()
    out = Path("docs/movate-cli-demo-deck.pptx")
    prs.save(out)
    print(f"✓ wrote {out} ({len(prs.slides)} slides)")
